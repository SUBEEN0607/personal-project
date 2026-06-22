"""
PE/VC 분기 보고 자동화 PoC — 발표 장표 PPTX 생성
초록 인포그래픽 스타일 · 화살표 · 시각화 중심
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ── 초록 컬러 팔레트 ──────────────────────────────
D_GREEN  = RGBColor(0x1b, 0x5e, 0x20)  # 다크 그린
M_GREEN  = RGBColor(0x2e, 0x7d, 0x32)  # 미디엄 그린
GREEN    = RGBColor(0x43, 0xa0, 0x47)  # 그린
L_GREEN  = RGBColor(0x81, 0xc7, 0x84)  # 라이트 그린
P_GREEN  = RGBColor(0xc8, 0xe6, 0xc9)  # 페일 그린
XP_GREEN = RGBColor(0xe8, 0xf5, 0xe9)  # 초연한 그린

BLACK    = RGBColor(0x1a, 0x1a, 0x1a)
D_GREY   = RGBColor(0x55, 0x55, 0x55)
GREY     = RGBColor(0x99, 0x99, 0x99)
L_GREY   = RGBColor(0xCC, 0xCC, 0xCC)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
BG       = RGBColor(0xFA, 0xFA, 0xF8)
BORDER   = RGBColor(0xE8, 0xE5, 0xE1)
WARM_BG  = RGBColor(0xF5, 0xF3, 0xEF)

FONT = "맑은 고딕"
W = Inches(13.333)
H = Inches(7.5)


# ── 헬퍼 함수들 ──────────────────────────────────

def _bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def _rect(slide, l, t, w, h, color=D_GREEN, border=None):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    r.fill.solid()
    r.fill.fore_color.rgb = color
    if border:
        r.line.color.rgb = border
        r.line.width = Pt(1)
    else:
        r.line.fill.background()
    return r

def _rounded(slide, l, t, w, h, color=WHITE, border=BORDER):
    r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    r.fill.solid()
    r.fill.fore_color.rgb = color
    r.line.color.rgb = border
    r.line.width = Pt(1)
    r.adjustments[0] = 0.05
    return r

def _circle(slide, l, t, size, color=D_GREEN):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, size, size)
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.line.fill.background()
    return c

def _text(slide, l, t, w, h, txt, sz=14, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT
    p.alignment = align
    p.space_after = Pt(0)
    return tb

def _multi(slide, l, t, w, h, lines, sz=12, color=BLACK, bold=False, spacing=1.4, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = FONT
        p.alignment = align
        p.space_after = Pt(int(spacing * 3))
    return tb

def _circle_text(slide, l, t, size, txt, bg_color=D_GREEN, txt_color=WHITE, txt_size=14):
    _circle(slide, l, t, size, bg_color)
    _text(slide, l, t, size, size, txt, sz=txt_size, color=txt_color, bold=True, align=PP_ALIGN.CENTER)

def _arrow_right(slide, l, t, w=Inches(0.5), h=Inches(0.3)):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    a.fill.solid()
    a.fill.fore_color.rgb = L_GREEN
    a.line.fill.background()
    return a

def _chevron(slide, l, t, w, h, color=M_GREEN):
    c = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, l, t, w, h)
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.line.fill.background()
    return c

def _pill(slide, l, t, w, h, txt, bg=P_GREEN, fg=D_GREEN, sz=10):
    r = _rounded(slide, l, t, w, h, color=bg, border=bg)
    _text(slide, l, t, w, h, txt, sz=sz, color=fg, bold=True, align=PP_ALIGN.CENTER)

def _progress_bar(slide, l, t, w, h, pct, bg_color=P_GREEN, fg_color=D_GREEN):
    _rounded(slide, l, t, w, h, color=bg_color, border=bg_color)
    fw = int(w * pct)
    if fw > 0:
        _rect(slide, l, t, fw, h, color=fg_color)

def _header(slide, label, title):
    _text(slide, Inches(0.8), Inches(0.35), Inches(4), Inches(0.25),
          label, sz=9, color=GREY, bold=True)
    _text(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.55),
          title, sz=28, color=BLACK, bold=True)
    _rect(slide, Inches(0.8), Inches(1.2), Inches(0.6), Pt(3), D_GREEN)
    _rect(slide, Inches(1.5), Inches(1.2), Inches(11), Pt(1), BORDER)

def _page(slide, n, total):
    _text(slide, Inches(12.0), Inches(7.0), Inches(1.2), Inches(0.3),
          f"{n} / {total}", sz=8, color=L_GREY, align=PP_ALIGN.RIGHT)

def _section(prs, title, subtitle=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, D_GREEN)
    _rect(s, Inches(0), Inches(0), W, Inches(0.06), L_GREEN)
    _rect(s, Inches(0), Inches(7.44), W, Inches(0.06), L_GREEN)
    _text(s, Inches(1.5), Inches(2.6), Inches(10), Inches(1),
          title, sz=40, color=WHITE, bold=True)
    if subtitle:
        _text(s, Inches(1.5), Inches(3.7), Inches(10), Inches(0.5),
              subtitle, sz=15, color=L_GREEN)

def _card(slide, l, t, w, h, icon_txt, title, lines, icon_bg=D_GREEN):
    _rounded(slide, l, t, w, h, WHITE, BORDER)
    _circle_text(slide, l + Inches(0.2), t + Inches(0.2),
                 Inches(0.4), icon_txt, bg_color=icon_bg, txt_size=11)
    _text(slide, l + Inches(0.75), t + Inches(0.22),
          w - Inches(1), Inches(0.3), title, sz=13, color=BLACK, bold=True)
    _multi(slide, l + Inches(0.2), t + Inches(0.7),
           w - Inches(0.4), h - Inches(0.8), lines, sz=10, color=D_GREY)


# ══════════════════════════════════════════════════
def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    TOTAL = 20

    # ═══════════════════════════════════════════════
    # 1. 표지
    # ═══════════════════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, D_GREEN)
    # 장식 원
    _circle(s, Inches(9.5), Inches(-1), Inches(5), M_GREEN)
    _circle(s, Inches(10.5), Inches(0), Inches(3), GREEN)
    _circle(s, Inches(-1), Inches(5.5), Inches(3), M_GREEN)
    # 상단 얇은 라인
    _rect(s, Inches(0), Inches(0), W, Inches(0.05), L_GREEN)
    _rect(s, Inches(0), Inches(7.45), W, Inches(0.05), L_GREEN)

    _text(s, Inches(1.2), Inches(1.5), Inches(6), Inches(0.3),
          "SDIC  ·  SKKU DIGITAL IT CONSULTING", sz=10, color=L_GREEN, bold=True)

    _rect(s, Inches(1.2), Inches(2.1), Inches(0.8), Pt(3), L_GREEN)

    _text(s, Inches(1.2), Inches(2.4), Inches(8), Inches(0.7),
          "PE/VC", sz=64, color=WHITE, bold=True)
    _text(s, Inches(1.2), Inches(3.4), Inches(8), Inches(0.5),
          "분기 보고 자동화 PoC", sz=26, color=P_GREEN)

    _text(s, Inches(1.2), Inches(4.4), Inches(8), Inches(0.4),
          "PoC 기획부터 구현까지 전 과정", sz=14, color=L_GREEN)

    _rect(s, Inches(1.2), Inches(5.3), Inches(2), Pt(1), L_GREEN)
    _text(s, Inches(1.2), Inches(5.6), Inches(4), Inches(0.3),
          "이수빈  |  2026.06", sz=13, color=L_GREEN)

    # ═══════════════════════════════════════════════
    # 2. 목차
    # ═══════════════════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "TABLE OF CONTENTS", "목차")

    sections = [
        ("01", "Problem & Motivation", "문제 정의 및 자동화 기회", D_GREEN),
        ("02", "Solution Design", "시스템 아키텍처 · 기술 스택 · 핵심 기능", M_GREEN),
        ("03", "Implementation", "주차별 구현 과정 — 실제 커밋 기반", GREEN),
        ("04", "Demo & Results", "성과 지표 · 앱 구조 시각화", L_GREEN),
        ("05", "Feedback & Next Steps", "피드백 대응 · 향후 로드맵", P_GREEN),
    ]
    for i, (num, title, desc, clr) in enumerate(sections):
        y = Inches(1.7) + Inches(i * 1.05)
        _circle_text(s, Inches(1.0), y, Inches(0.5), num, bg_color=clr, txt_size=13)
        _text(s, Inches(1.7), y + Inches(0.02), Inches(4), Inches(0.35),
              title, sz=18, color=BLACK, bold=True)
        _text(s, Inches(1.7), y + Inches(0.4), Inches(8), Inches(0.25),
              desc, sz=11, color=GREY)
    _page(s, 2, TOTAL)

    # ═══════════════════════════════════════════════
    # PART 1
    # ═══════════════════════════════════════════════
    _section(prs, "01  Problem & Motivation", "문제 정의 및 자동화 기회")

    # 3. Pain Point — 3카드 + 인포그래픽 숫자
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "PROBLEM", "GP가 LP 보고서를 만들 때의 현실")

    # 큰 숫자 인포그래픽
    stats = [
        ("4~8h", "보고서 1회\n작성 시간", D_GREEN),
        ("5+", "수동 계산\n지표 수", M_GREEN),
        ("3+", "수동 조회\n데이터 소스", GREEN),
    ]
    for i, (num, label, clr) in enumerate(stats):
        x = Inches(0.8) + Inches(i * 2.3)
        _rounded(s, x, Inches(1.6), Inches(2.0), Inches(1.8), XP_GREEN, P_GREEN)
        _text(s, x, Inches(1.7), Inches(2.0), Inches(0.8),
              num, sz=36, color=clr, bold=True, align=PP_ALIGN.CENTER)
        _text(s, x, Inches(2.5), Inches(2.0), Inches(0.7),
              label, sz=10, color=D_GREY, align=PP_ALIGN.CENTER)

    # 문제 카드 3개
    problems = [
        ("01", "엑셀 수작업",
         ["MOIC·IRR·DPI·RVPI·TVPI 포트폴리오사별 수동 계산",
          "공식 오류 빈번 → 보고서 신뢰도 하락",
          "반복 작업에 핵심 인력 시간 소모"]),
        ("02", "데이터 산재",
         ["DART 공시 → 웹에서 하나하나 수동 검색",
          "기준금리/환율 → 한국은행 사이트 수동 조회",
          "포트폴리오사 현재가치 → 감정평가 의존 (수주~수개월)"]),
        ("03", "분석 한계",
         ["J-Curve 시각화 → 엑셀 차트로 구현 어려움",
          "Exit 시나리오 IRR Sensitivity → 100개 셀 수작업",
          "Waterfall GP/LP 분배 → 복잡한 조건부 수식"]),
    ]
    for i, (num, title, lines) in enumerate(problems):
        x = Inches(0.8) + Inches(i * 4.1)
        _card(s, x, Inches(3.8), Inches(3.8), Inches(3.0), num, title, lines)
    _page(s, 4, TOTAL)

    # 4. Opportunity — API 소스 시각화
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "OPPORTUNITY", "공공 API + AI로 자동화 가능한 영역")

    apis = [
        ("D", "DART OpenAPI", "상장사 재무제표\n매출·영업이익·순이익", "0.95", D_GREEN),
        ("E", "ECOS 한국은행", "기준금리 월별 추이\n원/달러 환율", "0.90", M_GREEN),
        ("K", "KVIC 한국벤처투자", "모태펀드 분야별 현황\n연도별 결성 규모", "0.85", GREEN),
        ("C", "Claude AI", "LP 코멘터리 자동 생성\n자연어 질문 답변(RAG)", "1.0", D_GREEN),
        ("N", "네이버 금융/뉴스", "상장사 시가총액 크롤링\n뉴스 모니터링", "0.80", L_GREEN),
    ]
    for i, (icon, name, desc, reliability, clr) in enumerate(apis):
        y = Inches(1.6) + Inches(i * 1.1)
        # 아이콘 원
        _circle_text(s, Inches(0.8), y + Inches(0.05), Inches(0.45), icon, bg_color=clr, txt_size=14)
        # 이름
        _text(s, Inches(1.5), y, Inches(2.5), Inches(0.3),
              name, sz=14, color=BLACK, bold=True)
        # 설명
        _text(s, Inches(1.5), y + Inches(0.35), Inches(3.5), Inches(0.6),
              desc, sz=10, color=D_GREY)
        # 프로그레스 바
        pct = float(reliability)
        bar_l = Inches(5.5)
        bar_w = Inches(6.5)
        bar_h = Inches(0.2)
        _rounded(s, bar_l, y + Inches(0.2), bar_w, bar_h, P_GREEN, P_GREEN)
        fill_w = int(bar_w * pct)
        _rect(s, bar_l, y + Inches(0.2), fill_w, bar_h, clr)
        _text(s, bar_l + bar_w + Inches(0.1), y + Inches(0.12), Inches(0.5), Inches(0.3),
              f"{int(pct*100)}%", sz=10, color=clr, bold=True)
    _page(s, 5, TOTAL)

    # ═══════════════════════════════════════════════
    # PART 2
    # ═══════════════════════════════════════════════
    _section(prs, "02  Solution Design", "아키텍처 · 기술 스택 · 핵심 기능")

    # 5. 아키텍처 — 화살표 플로우
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "ARCHITECTURE", "시스템 아키텍처 — 데이터 파이프라인")

    steps = [
        ("01", "INPUT", ["CSV/Excel 업로드", "자동 밸류에이션", "(네이버 시총 크롤링)"], D_GREEN),
        ("02", "CALCULATE", ["calculator.py", "irr.py / simulator.py", "MOIC·IRR·DPI·RVPI·TVPI"], M_GREEN),
        ("03", "ENRICH", ["DART 재무 조회", "ECOS 거시지표", "KVIC 벤치마크"], GREEN),
        ("04", "ANALYZE", ["Claude AI 코멘터리", "RAG 자연어 질문", "뉴스 모니터링"], M_GREEN),
        ("05", "OUTPUT", ["Streamlit 대시보드", "PDF 보고서", "SQLite 분기 저장"], D_GREEN),
    ]
    for i, (num, label, items, clr) in enumerate(steps):
        x = Inches(0.3) + Inches(i * 2.6)
        # 상단 원 번호
        _circle_text(s, x + Inches(0.8), Inches(1.6), Inches(0.5), num, bg_color=clr, txt_size=12)
        # 라벨
        _text(s, x, Inches(2.2), Inches(2.3), Inches(0.3),
              label, sz=11, color=clr, bold=True, align=PP_ALIGN.CENTER)
        # 카드
        _rounded(s, x, Inches(2.6), Inches(2.3), Inches(3.0), WHITE, clr)
        _multi(s, x + Inches(0.2), Inches(2.8), Inches(1.9), Inches(2.5),
               items, sz=11, color=BLACK, spacing=2.0)
        # 화살표
        if i < 4:
            _arrow_right(s, x + Inches(2.3), Inches(3.8), Inches(0.35), Inches(0.25))

    # 하단 범례
    _text(s, Inches(0.8), Inches(6.2), Inches(10), Inches(0.3),
          "ThreadPoolExecutor(6) 병렬 처리  ·  @st.cache_data 캐싱  ·  SQLite 스냅샷 저장",
          sz=10, color=GREY)
    _page(s, 7, TOTAL)

    # 6. 기술 스택 — 카드 그리드
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "TECH STACK", "기술 스택")

    stacks = [
        ("F", "Frontend", ["Streamlit", "Plotly 인터랙티브 차트", "Pretendard 타이포그래피", "CSS 디자인 시스템"], D_GREEN),
        ("B", "Backend", ["Python 3.11", "pandas · scipy", "sqlite3", "concurrent.futures"], M_GREEN),
        ("A", "APIs", ["DART OpenAPI", "ECOS (한국은행)", "KVIC (한국벤처투자)", "네이버 금융/뉴스"], GREEN),
        ("I", "AI", ["Claude claude-haiku-4-5", "sentence-transformers", "BeautifulSoup", "RAG 파이프라인"], M_GREEN),
        ("O", "Output", ["fpdf2 (PDF)", "python-pptx (PPTX)", "SQLite (분기 저장)", "st.download_button"], D_GREEN),
    ]
    for i, (icon, title, items, clr) in enumerate(stacks):
        x = Inches(0.3) + Inches(i * 2.6)
        _rounded(s, x, Inches(1.6), Inches(2.3), Inches(4.5), WHITE, BORDER)
        _circle_text(s, x + Inches(0.85), Inches(1.8), Inches(0.5), icon, bg_color=clr, txt_size=14)
        _text(s, x, Inches(2.4), Inches(2.3), Inches(0.3),
              title, sz=14, color=BLACK, bold=True, align=PP_ALIGN.CENTER)
        _rect(s, x + Inches(0.5), Inches(2.85), Inches(1.3), Pt(2), clr)
        for j, item in enumerate(items):
            _pill(s, x + Inches(0.15), Inches(3.1) + Inches(j * 0.55),
                  Inches(2.0), Inches(0.4), item, bg=XP_GREEN, fg=D_GREEN, sz=10)
    _page(s, 8, TOTAL)

    # 7. 핵심 기능 — vs 엑셀
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "FEATURES", "핵심 기능 — 엑셀 vs 자동화")

    features = [
        ("01", "MOIC·IRR\n자동 계산", "CSV 즉시 XIRR 기반\n8개사 2초 이내", "수동 XIRR\n공식 입력"),
        ("02", "IRR Sensitivity\nMatrix", "Exit배수×보유기간\n100셀 히트맵", "매번\n수동 계산"),
        ("03", "J-Curve\n시각화", "누적 현금흐름 차트\n손익분기 자동 표시", "차트\n수작업"),
        ("04", "Waterfall\n분배", "Hurdle/Carry 4단계\n실시간 시뮬레이션", "복잡한\n조건부 수식"),
        ("05", "AI\n코멘터리", "Claude LP 보고서\nRAG 질문 답변", "불가능"),
        ("06", "자동\n밸류에이션", "네이버 시총 크롤링\n섹터 P/S 배수", "불가능"),
    ]
    for i, (num, title, ours, theirs) in enumerate(features):
        col = i % 3
        row = i // 3
        x = Inches(0.5) + Inches(col * 4.2)
        y = Inches(1.6) + Inches(row * 2.8)

        _rounded(s, x, y, Inches(3.9), Inches(2.5), WHITE, BORDER)
        _circle_text(s, x + Inches(0.15), y + Inches(0.15), Inches(0.35), num, bg_color=D_GREEN, txt_size=9)
        _text(s, x + Inches(0.6), y + Inches(0.15), Inches(2), Inches(0.5),
              title, sz=12, color=BLACK, bold=True)

        # 자동화 (초록)
        _pill(s, x + Inches(0.15), y + Inches(0.85), Inches(0.6), Inches(0.22),
              "자동화", bg=D_GREEN, fg=WHITE, sz=8)
        _text(s, x + Inches(0.85), y + Inches(0.8), Inches(2.8), Inches(0.6),
              ours, sz=10, color=BLACK)

        # 엑셀 (회색)
        _pill(s, x + Inches(0.15), y + Inches(1.65), Inches(0.6), Inches(0.22),
              "엑셀", bg=BORDER, fg=GREY, sz=8)
        _text(s, x + Inches(0.85), y + Inches(1.6), Inches(2.8), Inches(0.5),
              theirs, sz=10, color=GREY)
    _page(s, 9, TOTAL)

    # ═══════════════════════════════════════════════
    # PART 3: Implementation
    # ═══════════════════════════════════════════════
    _section(prs, "03  Implementation", "주차별 구현 과정 — 실제 git 커밋 기반")

    # 8. 타임라인 개요
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "TIMELINE", "4주간 구현 타임라인")

    weeks = [
        ("W1", "5/21~25", "PoC 스켈레톤", "GitHub repo · calculator.py\ncommentary.py · app.py", D_GREEN, 1.0),
        ("W2", "5/25~26", "API 연동 + UI", "DART·ECOS·KVIC 연동\n탭별 PDF · Plotly 차트", M_GREEN, 1.0),
        ("W3", "5/27", "고급 분석", "Waterfall · IRR Sensitivity\n뉴스 모니터링 · KSD API", GREEN, 1.0),
        ("W4", "6/22", "자동화 + 리디자인", "자동 밸류에이션\nUI 전면 재설계", M_GREEN, 0.8),
    ]

    # 타임라인 바
    bar_y = Inches(2.0)
    _rect(s, Inches(1.5), bar_y + Inches(0.18), Inches(10.5), Pt(4), P_GREEN)

    for i, (wk, dates, title, desc, clr, pct) in enumerate(weeks):
        x = Inches(1.5) + Inches(i * 2.8)
        # 원 위 점
        _circle(s, x + Inches(0.15), bar_y, Inches(0.4), clr)
        _text(s, x + Inches(0.15), bar_y, Inches(0.4), Inches(0.4),
              wk, sz=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        # 날짜
        _text(s, x - Inches(0.3), bar_y + Inches(0.5), Inches(1.3), Inches(0.2),
              dates, sz=9, color=GREY, align=PP_ALIGN.CENTER)
        # 카드
        _rounded(s, x - Inches(0.5), bar_y + Inches(0.9), Inches(2.5), Inches(3.5), WHITE, clr)
        _text(s, x - Inches(0.3), bar_y + Inches(1.05), Inches(2.1), Inches(0.3),
              title, sz=13, color=clr, bold=True)
        _multi(s, x - Inches(0.3), bar_y + Inches(1.5), Inches(2.1), Inches(2.5),
               desc.split("\n"), sz=11, color=BLACK, spacing=1.5)
        # 진행률 바
        _progress_bar(s, x - Inches(0.3), bar_y + Inches(3.8), Inches(2.0), Inches(0.15), pct)
    _page(s, 11, TOTAL)

    # 9. Week 1 상세
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "WEEK 1  ·  5/21 ~ 5/25", "PoC 스켈레톤 — 기본 구조 완성")
    # 좌: 구현 내용
    _rounded(s, Inches(0.8), Inches(1.6), Inches(6.0), Inches(5.2), WHITE, P_GREEN)
    _text(s, Inches(1.0), Inches(1.7), Inches(3), Inches(0.3),
          "구현 내용", sz=14, color=D_GREEN, bold=True)
    impl_lines = [
        "GitHub repo 생성 (SUBEEN0607/personal-project)",
        "CLAUDE.md 작성 — 프로젝트 컨벤션 정의",
        "",
        "calculator.py — MOIC/IRR/DPI/RVPI/TVPI 계산 로직",
        "commentary.py — Claude API (claude-haiku-4-5) 코멘터리",
        "app.py — Streamlit UI (CSV → 지표 → AI 코멘터리)",
        "sample_portfolio.csv — 8개사 가상 포트폴리오",
        "",
        ".gitignore, .env.example, requirements.txt",
        "로컬 환경 구성 + localhost:8501 실행 확인",
    ]
    _multi(s, Inches(1.0), Inches(2.1), Inches(5.5), Inches(4.5),
           impl_lines, sz=11, color=BLACK, spacing=1.5)

    # 우: 커밋 + 성과
    _rounded(s, Inches(7.2), Inches(1.6), Inches(5.3), Inches(2.8), XP_GREEN, P_GREEN)
    _text(s, Inches(7.4), Inches(1.7), Inches(3), Inches(0.3),
          "주요 커밋", sz=12, color=D_GREEN, bold=True)
    commits = [
        "f2ddd64  calculator.py - MOIC/IRR/TVPI",
        "6046565  commentary.py - Claude API",
        "863218f  app.py - Streamlit PoC UI",
        "f2ecb22  sample_portfolio.csv",
    ]
    _multi(s, Inches(7.4), Inches(2.1), Inches(5.0), Inches(2.0),
           commits, sz=9, color=D_GREY, spacing=1.5)

    _rounded(s, Inches(7.2), Inches(4.8), Inches(5.3), Inches(2.0), WHITE, D_GREEN)
    _text(s, Inches(7.4), Inches(4.9), Inches(3), Inches(0.3),
          "성과", sz=12, color=D_GREEN, bold=True)
    _multi(s, Inches(7.4), Inches(5.3), Inches(5.0), Inches(1.3),
           ["CSV → 5개 지표 자동 계산 → AI 코멘터리 1클릭",
            "엑셀 2시간 → 10초 (계산 부분 기준)"],
           sz=11, color=BLACK, spacing=1.5)
    _page(s, 12, TOTAL)

    # 10. Week 2 상세
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "WEEK 2  ·  5/25 ~ 5/26", "API 연동 + 대시보드 UI 고도화")

    _rounded(s, Inches(0.8), Inches(1.6), Inches(6.0), Inches(5.2), WHITE, P_GREEN)
    _text(s, Inches(1.0), Inches(1.7), Inches(3), Inches(0.3),
          "구현 내용", sz=14, color=D_GREEN, bold=True)
    _multi(s, Inches(1.0), Inches(2.1), Inches(5.5), Inches(4.5),
           ["Plotly 차트 3종 (MOIC 바차트, 산점도, 섹터 파이)",
            "rag.py — sentence-transformers 자연어 Q&A",
            "irr.py — XIRR 기반 IRR + J-Curve 데이터",
            "simulator.py — Exit 배수별 IRR 시뮬레이터",
            "db.py — SQLite 분기별 스냅샷 저장",
            "report.py — fpdf2 LP 보고서 PDF",
            "",
            "DART OpenAPI — 기업 검색 + 재무제표",
            "ECOS — 기준금리(M) + 환율(일별→월별)",
            "KVIC — 모태펀드 분야별/연도별 벤치마크",
            "탭별 AI 해석 PDF 5종"],
           sz=11, color=BLACK, spacing=1.4)

    _rounded(s, Inches(7.2), Inches(1.6), Inches(5.3), Inches(2.5), XP_GREEN, P_GREEN)
    _text(s, Inches(7.4), Inches(1.7), Inches(3), Inches(0.3),
          "해결한 이슈", sz=12, color=D_GREEN, bold=True)
    _multi(s, Inches(7.4), Inches(2.1), Inches(5.0), Inches(2.0),
           ["fpdf2 uni=True 제거 → 한글 PDF 정상",
            "ECOS cycle MM→M (기준금리 조회 실패 수정)",
            "dart-fss → OpenAPI 직접 호출 (속도 개선)",
            "KVIC 연도 필터 → Python측 처리"],
           sz=10, color=D_GREY, spacing=1.5)

    _rounded(s, Inches(7.2), Inches(4.5), Inches(5.3), Inches(2.3), WHITE, M_GREEN)
    _text(s, Inches(7.4), Inches(4.6), Inches(3), Inches(0.3),
          "커밋 하이라이트", sz=12, color=M_GREEN, bold=True)
    _multi(s, Inches(7.4), Inches(5.0), Inches(5.0), Inches(1.6),
           ["818ed5c  Plotly + 자연어 Q&A",
            "a1f6b15  KVIC API 연동",
            "786c377  DART OpenAPI 직접 호출",
            "7a265fc  탭별 AI PDF 보고서"],
           sz=9, color=D_GREY, spacing=1.5)
    _page(s, 13, TOTAL)

    # 11. Week 3
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "WEEK 3  ·  5/27", "고급 분석 + UI 고도화")

    _rounded(s, Inches(0.8), Inches(1.6), Inches(6.0), Inches(5.2), WHITE, P_GREEN)
    _text(s, Inches(1.0), Inches(1.7), Inches(3), Inches(0.3),
          "구현 내용", sz=14, color=D_GREEN, bold=True)
    _multi(s, Inches(1.0), Inches(2.1), Inches(5.5), Inches(4.5),
           ["Waterfall 분배 계산기 (Hurdle/Catch-up/Carry 4단계)",
            "IRR Sensitivity Matrix — 히트맵 (Exit배수 × 보유기간)",
            "포트폴리오사 뉴스 모니터링 (Naver API)",
            "KSD 오픈API 클라이언트 — 30개 엔드포인트",
            "",
            "탭 7개→5개 통합 (펀드추이·투자분석 병합)",
            "ILPA 스타일 대시보드 + 지표 계층화",
            "표지 랜딩 페이지 (배경 그라디언트 + SDIC 로고)",
            "API 로고 SVG 아이콘",
            "전체 차트 #2e7d32 초록 통일"],
           sz=11, color=BLACK, spacing=1.4)

    _rounded(s, Inches(7.2), Inches(1.6), Inches(5.3), Inches(2.2), XP_GREEN, P_GREEN)
    _text(s, Inches(7.4), Inches(1.7), Inches(3), Inches(0.3),
          "디자인 개선", sz=12, color=D_GREEN, bold=True)
    _multi(s, Inches(7.4), Inches(2.1), Inches(5.0), Inches(1.8),
           ["Pretendard 폰트 적용 (가독성 향상)",
            "호버 시 카드 부상 효과",
            "지표 용어 해설 Expander",
            "슬라이더 목표IRR 기준선 실시간 반영"],
           sz=10, color=D_GREY, spacing=1.5)

    _rounded(s, Inches(7.2), Inches(4.2), Inches(5.3), Inches(2.6), WHITE, GREEN)
    _text(s, Inches(7.4), Inches(4.3), Inches(3), Inches(0.3),
          "커밋 하이라이트", sz=12, color=GREEN, bold=True)
    _multi(s, Inches(7.4), Inches(4.7), Inches(5.0), Inches(2.0),
           ["a8f0638  Waterfall 분배 계산기",
            "c61106a  IRR Sensitivity + 뉴스 모니터링",
            "a525682  KSD 오픈API 클라이언트",
            "754ca7b  표지 페이지 + 트리맵",
            "e6e575d  피드백 반영 리팩토링"],
           sz=9, color=D_GREY, spacing=1.5)
    _page(s, 14, TOTAL)

    # 12. Week 4
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "WEEK 4  ·  6/22", "자동 밸류에이션 + UI 전면 재설계")

    # 좌: 자동 밸류에이션 플로우
    _rounded(s, Inches(0.8), Inches(1.6), Inches(6.0), Inches(2.6), WHITE, D_GREEN)
    _text(s, Inches(1.0), Inches(1.7), Inches(3), Inches(0.3),
          "자동 밸류에이션 파이프라인", sz=13, color=D_GREEN, bold=True)

    flow = [
        ("회사명\n입력", D_GREEN), ("DART\n검색", M_GREEN),
        ("상장?\n판단", GREEN), ("네이버\n시총", M_GREEN), ("현재가치\n계산", D_GREEN),
    ]
    for i, (label, clr) in enumerate(flow):
        x = Inches(1.0) + Inches(i * 1.9)
        _rounded(s, x, Inches(2.2), Inches(1.2), Inches(0.8), clr, clr)
        _text(s, x, Inches(2.2), Inches(1.2), Inches(0.8),
              label, sz=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        if i < 4:
            _arrow_right(s, x + Inches(1.2), Inches(2.45), Inches(0.6), Inches(0.2))

    _text(s, Inches(1.0), Inches(3.3), Inches(5.5), Inches(0.5),
          "비상장사 fallback: DART 매출 × 섹터 P/S 배수 (20개 섹터 테이블)",
          sz=10, color=D_GREY)

    # 우: UI 리디자인
    _rounded(s, Inches(7.2), Inches(1.6), Inches(5.3), Inches(2.6), XP_GREEN, P_GREEN)
    _text(s, Inches(7.4), Inches(1.7), Inches(3), Inches(0.3),
          "UI 전면 재설계", sz=13, color=D_GREEN, bold=True)
    _multi(s, Inches(7.4), Inches(2.1), Inches(5.0), Inches(2.0),
           ["Apple-inspired 미니멀 디자인 시스템",
            "CSS 전면 교체 (warm neutral + green accent)",
            "커버 페이지 리디자인 (대형 타이포)",
            "플랫 카드 (그림자 제거, 얇은 보더)",
            "PDF 보고서 레이아웃 전면 재설계"],
           sz=10, color=BLACK, spacing=1.5)

    # 하단: 피드백 대응
    _rounded(s, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.2), WHITE, BORDER)
    _text(s, Inches(1.0), Inches(4.7), Inches(3), Inches(0.3),
          "피드백 5건 대응", sz=13, color=D_GREEN, bold=True)
    fb_items = [
        ("Q1", "PE/VC KPI 차이 → 설계 완료"),
        ("Q2", "Waterfall 위치 → 적절 (Tab3 유지)"),
        ("Q3", "KVIC API → 이미 구현됨"),
        ("Q4", "PDF→PPTX → python-pptx 추가"),
        ("Q5", "데이터 역추출 → 자동 밸류에이션 완료"),
    ]
    for i, (q, txt) in enumerate(fb_items):
        x = Inches(1.0) + Inches((i % 3) * 3.9)
        y = Inches(5.15) + Inches((i // 3) * 0.65)
        _circle_text(s, x, y, Inches(0.3), q, bg_color=GREEN, txt_size=8)
        _text(s, x + Inches(0.4), y + Inches(0.02), Inches(3.3), Inches(0.25),
              txt, sz=10, color=BLACK)
    _page(s, 15, TOTAL)

    # ═══════════════════════════════════════════════
    # PART 4
    # ═══════════════════════════════════════════════
    _section(prs, "04  Demo & Results", "성과 지표 · 앱 구조")

    # 13. Before vs After
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "RESULTS", "Before vs After")

    metrics = [
        ("보고서 작성 시간", "4~8시간", "~30분", "87%↓"),
        ("수동 계산 항목", "5지표 × N사", "CSV 즉시", "100%"),
        ("분석 범위", "MOIC·IRR만", "J-Curve, Waterfall\nSensitivity, AI", "6배↑"),
        ("데이터 소스", "웹 수동 검색", "5개 API 자동", "실시간"),
        ("밸류에이션", "감정평가 수주", "네이버 크롤링\n<10초", "즉시"),
    ]
    # 컬럼 헤더
    headers = ["", "Before", "After", "개선"]
    hx = [Inches(0.8), Inches(4.0), Inches(7.5), Inches(11.0)]
    for j, (hdr, x) in enumerate(zip(headers, hx)):
        bg_c = BORDER if j == 0 else (RGBColor(0xEE,0xEE,0xEE) if j == 1 else (XP_GREEN if j == 2 else P_GREEN))
        _pill(s, x, Inches(1.6), Inches(1.8) if j > 0 else Inches(2.8), Inches(0.3),
              hdr, bg=bg_c, fg=D_GREEN if j >= 2 else GREY, sz=10)

    for i, (item, before, after, improve) in enumerate(metrics):
        y = Inches(2.2) + Inches(i * 1.0)
        # 항목
        _text(s, Inches(0.8), y, Inches(2.8), Inches(0.3),
              item, sz=13, color=BLACK, bold=True)
        # Before (회색 박스)
        _rounded(s, Inches(4.0), y, Inches(3.0), Inches(0.7), WARM_BG, BORDER)
        _text(s, Inches(4.1), y + Inches(0.05), Inches(2.8), Inches(0.6),
              before, sz=11, color=GREY)
        # After (초록 박스)
        _rounded(s, Inches(7.5), y, Inches(3.0), Inches(0.7), XP_GREEN, P_GREEN)
        _text(s, Inches(7.6), y + Inches(0.05), Inches(2.8), Inches(0.6),
              after, sz=11, color=D_GREEN, bold=True)
        # 개선 (원)
        _circle_text(s, Inches(11.0), y + Inches(0.05), Inches(0.6), improve,
                     bg_color=D_GREEN, txt_size=10)
    _page(s, 17, TOTAL)

    # 14. 앱 구조
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "APP STRUCTURE", "앱 구조 — 5탭 · 12개 파일")

    tabs = [
        ("01", "대시보드", "MOIC·IRR Hero\nDPI/RVPI/TVPI 보조\n4종 차트\nPDF 보고서", D_GREEN),
        ("02", "펀드 추이", "J-Curve\n분기별 추이\nAI 해석 PDF", M_GREEN),
        ("03", "투자 분석", "DART 재무 조회\n시나리오 시뮬레이터\nIRR Sensitivity\nWaterfall", GREEN),
        ("04", "시장 벤치마크", "ECOS 기준금리\n환율 추이\nKVIC 모태펀드", L_GREEN),
        ("05", "AI 분석", "Claude 코멘터리\n자연어 질문(RAG)\n뉴스 모니터링", M_GREEN),
    ]
    for i, (num, title, desc, clr) in enumerate(tabs):
        x = Inches(0.3) + Inches(i * 2.6)
        _rounded(s, x, Inches(1.6), Inches(2.3), Inches(5.0), WHITE, clr)
        _circle_text(s, x + Inches(0.85), Inches(1.8), Inches(0.5), num, bg_color=clr, txt_size=12)
        _text(s, x, Inches(2.4), Inches(2.3), Inches(0.3),
              title, sz=14, color=BLACK, bold=True, align=PP_ALIGN.CENTER)
        _rect(s, x + Inches(0.5), Inches(2.85), Inches(1.3), Pt(2), clr)
        _multi(s, x + Inches(0.2), Inches(3.1), Inches(1.9), Inches(3.3),
               desc.split("\n"), sz=11, color=D_GREY, spacing=2.0)
    _page(s, 18, TOTAL)

    # ═══════════════════════════════════════════════
    # PART 5
    # ═══════════════════════════════════════════════
    _section(prs, "05  Feedback & Next Steps", "피드백 대응 · 향후 로드맵")

    # 15. 향후 로드맵 — 타임라인
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "ROADMAP", "향후 로드맵")

    phases = [
        ("단기", "7월", [
            "PE/VC별 KPI 조건부 표시",
            "PPTX 보고서 자동 생성",
            "PME 벤치마크 (pykrx)",
        ], D_GREEN),
        ("중기", "8~9월", [
            "Streamlit Cloud 배포",
            "LP별 개인화 보고서",
            "포트폴리오 모니터링 대시보드",
        ], M_GREEN),
        ("장기", "10월~", [
            "비상장사 밸류에이션 고도화",
            "Excel 템플릿 자동 생성",
            "멀티 펀드 관리",
        ], GREEN),
    ]
    for i, (phase, period, items, clr) in enumerate(phases):
        x = Inches(0.8) + Inches(i * 4.1)
        # 상단 페이즈 라벨
        _pill(s, x, Inches(1.6), Inches(1.2), Inches(0.35), phase, bg=clr, fg=WHITE, sz=12)
        _text(s, x + Inches(1.4), Inches(1.6), Inches(1.5), Inches(0.35),
              period, sz=13, color=clr, bold=True)

        _rounded(s, x, Inches(2.2), Inches(3.7), Inches(4.0), WHITE, clr)
        for j, item in enumerate(items):
            y = Inches(2.5) + Inches(j * 1.0)
            _circle_text(s, x + Inches(0.2), y, Inches(0.35), str(j+1),
                         bg_color=clr, txt_size=10)
            _text(s, x + Inches(0.7), y + Inches(0.02), Inches(2.8), Inches(0.3),
                  item, sz=12, color=BLACK)
    _page(s, 20, TOTAL)

    # 16. Q&A
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, D_GREEN)
    _circle(s, Inches(9), Inches(-1.5), Inches(6), M_GREEN)
    _circle(s, Inches(10), Inches(-0.5), Inches(4), GREEN)
    _circle(s, Inches(-1.5), Inches(5), Inches(4), M_GREEN)
    _rect(s, Inches(0), Inches(0), W, Inches(0.05), L_GREEN)
    _rect(s, Inches(0), Inches(7.45), W, Inches(0.05), L_GREEN)

    _text(s, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
          "Q&A", sz=72, color=WHITE, bold=True)
    _rect(s, Inches(1.5), Inches(3.5), Inches(2), Pt(2), L_GREEN)
    _text(s, Inches(1.5), Inches(3.8), Inches(10), Inches(0.4),
          "이수빈  ·  SDIC  ·  PE/VC 분기 보고 자동화 PoC", sz=14, color=L_GREEN)
    _text(s, Inches(1.5), Inches(4.4), Inches(10), Inches(0.3),
          "GitHub: SUBEEN0607/personal-project", sz=12, color=P_GREEN)

    # ── 저장 ──
    path = "PE_VC_PoC_Presentation.pptx"
    prs.save(path)
    print(f"Done: {path}")
    return path


if __name__ == "__main__":
    build()
