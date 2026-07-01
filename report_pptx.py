"""
LP 보고서 PPTX — 컨설팅 장표 수준 재작성
모든 좌표를 픽셀 단위로 설계하여 겹침/공백 완전 제거
"""
import io, os, tempfile, datetime
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 팔레트 ──────────────────────────────────────────
C_PRIMARY   = RGBColor(0x1B, 0x5E, 0x20)
C_SECONDARY = RGBColor(0x2E, 0x7D, 0x32)
C_ACCENT    = RGBColor(0x43, 0xA0, 0x47)
C_LIGHT     = RGBColor(0xA5, 0xD6, 0xA7)
C_PALE      = RGBColor(0xC8, 0xE6, 0xC9)
C_XPALE     = RGBColor(0xE8, 0xF5, 0xE9)
C_BG        = RGBColor(0xFA, 0xFA, 0xF8)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK     = RGBColor(0x1A, 0x1A, 0x1A)
C_DARK      = RGBColor(0x33, 0x33, 0x33)
C_GREY      = RGBColor(0x77, 0x77, 0x77)
C_LGREY     = RGBColor(0xBB, 0xBB, 0xBB)
C_BORDER    = RGBColor(0xE0, 0xE0, 0xE0)
C_RED       = RGBColor(0xC6, 0x28, 0x28)
C_ORANGE    = RGBColor(0xFF, 0x98, 0x00)
C_BEIGE     = RGBColor(0xF5, 0xF0, 0xE8)
C_INSIGHT   = RGBColor(0xF7, 0xF5, 0xF0)

FONT = "맑은 고딕"
SW, SH = Inches(13.333), Inches(7.5)

# 레이아웃 상수 — 모든 슬라이드 공통
M_LEFT   = Inches(0.6)          # 좌측 마진
M_RIGHT  = Inches(12.73)        # 우측 끝 (13.333-0.6)
C_WIDTH  = Inches(12.13)        # 콘텐츠 폭
TOP_BAR  = Inches(0.06)         # 상단 바 높이
HDR_Y    = Inches(0.10)         # 제목 Y
MSG_Y    = Inches(0.88)         # 핵심 메시지 Y
BODY_Y   = Inches(1.85)         # 본문 시작 Y (헤더+메시지 아래 충분한 간격)
BOT_Y    = Inches(7.10)         # 페이지 번호 Y

# ── 기본 헬퍼 ───────────────────────────────────────
def _bg(slide, color=C_BG):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def _shape(s, l, t, w, h, fill, border=None, radius=False):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    r = s.shapes.add_shape(kind, l, t, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = fill
    if border:
        r.line.color.rgb = border; r.line.width = Pt(0.75)
    else:
        r.line.fill.background()
    if radius:
        r.adjustments[0] = 0.03
    return r

def _txt(s, l, t, w, h, text, sz=10, color=C_BLACK, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = str(text)
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold
    p.font.name = font; p.alignment = align; p.space_after = Pt(0)
    tf.auto_size = None
    return tb

def _multiline(s, l, t, w, h, lines, sz=10, color=C_DARK, spacing=4, bold=False):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(line)
        p.font.size = Pt(sz); p.font.color.rgb = color
        p.font.name = FONT; p.space_after = Pt(spacing); p.font.bold = bold
    return tb

# ── 컴포넌트 ───────────────────────────────────────
def _header(s, section_label, title, key_message=""):
    _txt(s, M_LEFT, HDR_Y, Inches(8), Inches(0.45), title,
         sz=26, color=C_BLACK, bold=True)
    _txt(s, Inches(9.5), HDR_Y + Inches(0.10), Inches(3.2), Inches(0.35),
         section_label, sz=10, color=C_GREY, align=PP_ALIGN.RIGHT)
    # 레퍼런스 스타일: 두꺼운 초록 언더바
    _shape(s, M_LEFT, Inches(0.70), C_WIDTH, Pt(3.5), C_PRIMARY)
    if key_message:
        _txt(s, M_LEFT, Inches(0.92), C_WIDTH, Inches(0.32),
             key_message, sz=13, color=C_BLACK, bold=True)

def _page_num(s, n, total):
    _txt(s, Inches(12.0), BOT_Y, Inches(1.2), Inches(0.2),
         f"{n} / {total}", sz=7, color=C_LGREY, align=PP_ALIGN.RIGHT)

def _kpi_card(s, l, t, w, h, label, value, sub="", val_color=C_PRIMARY):
    _shape(s, l, t, w, h, C_WHITE, C_BORDER, radius=True)
    _shape(s, l, t, w, Pt(3.5), val_color)   # 상단 컬러 액센트 스트라이프
    _txt(s, l, t + Inches(0.14), w, Inches(0.15), label,
         sz=8, color=C_GREY, bold=True, align=PP_ALIGN.CENTER)
    vy = t + Inches(0.32) if h >= Inches(1.0) else t + Inches(0.26)
    _txt(s, l, vy, w, Inches(0.35), str(value),
         sz=22 if h >= Inches(1.0) else 18, color=val_color, bold=True, align=PP_ALIGN.CENTER)
    if sub:
        _txt(s, l, t + h - Inches(0.22), w, Inches(0.18), sub,
             sz=7, color=C_GREY, align=PP_ALIGN.CENTER)

def _formula_box(s, l, t, w, h, title, formula, example, criteria):
    """레퍼런스 Appendix 스타일 — 공식+예시+해석기준 박스 (콘텐츠에 맞춘 타이트한 배치)"""
    _shape(s, l, t, w, h, C_WHITE, C_BORDER, radius=True)
    _txt(s, l + Inches(0.14), t + Inches(0.12), w - Inches(0.28), Inches(0.24),
         title, sz=10.5, color=C_BLACK, bold=True)
    _shape(s, l + Inches(0.14), t + Inches(0.42), w - Inches(0.28), Inches(0.36), C_PALE, radius=True)
    _txt(s, l + Inches(0.20), t + Inches(0.46), w - Inches(0.40), Inches(0.30),
         formula, sz=9, color=C_PRIMARY, bold=True)
    _txt(s, l + Inches(0.14), t + Inches(0.86), w - Inches(0.28), Inches(0.42),
         example, sz=8, color=C_DARK)
    crit_y = t + Inches(1.32)
    crit_h = max(h - Inches(1.32) - Inches(0.10), Inches(0.40))
    _shape(s, l + Inches(0.14), crit_y, w - Inches(0.28), crit_h, C_BEIGE, radius=True)
    _txt(s, l + Inches(0.20), crit_y + Inches(0.06), w - Inches(0.40), crit_h - Inches(0.10),
         f"해석기준: {criteria}", sz=7.5, color=C_DARK)

def _process_box(s, l, t, w, h, num, title, desc, bg=C_XPALE, num_bg=C_PRIMARY):
    """레퍼런스 스타일 — 번호 아이콘 + 제목 + 설명 박스"""
    _shape(s, l, t, w, h, bg, C_BORDER, radius=True)
    _shape(s, l + Inches(0.12), t + Inches(0.12), Inches(0.30), Inches(0.30), num_bg, radius=True)
    _txt(s, l + Inches(0.12), t + Inches(0.14), Inches(0.30), Inches(0.26),
         str(num), sz=11, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    _txt(s, l + Inches(0.55), t + Inches(0.12), w - Inches(0.70), Inches(0.22),
         title, sz=10, color=C_BLACK, bold=True)
    _txt(s, l + Inches(0.12), t + Inches(0.46), w - Inches(0.24), h - Inches(0.55),
         desc, sz=8, color=C_DARK)

def _process_list(s, l, t, w, items, box_h=Inches(0.46), gap=Inches(0.08), title="ANALYSIS PROCESS"):
    """ANALYSIS 텍스트 블록을 번호 박스 리스트로 구조화"""
    _txt(s, l, t, w, Inches(0.18), title, sz=8.5, color=C_GREY, bold=True)
    y0 = t + Inches(0.34)
    for i, text in enumerate(items):
        y = y0 + i * (box_h + gap)
        _shape(s, l, y, w, box_h, C_WHITE, C_BORDER, radius=True)
        _shape(s, l + Inches(0.10), y + Inches(0.08), Inches(0.28), Inches(0.28), C_PRIMARY, radius=True)
        _txt(s, l + Inches(0.10), y + Inches(0.105), Inches(0.28), Inches(0.24),
             str(i + 1), sz=9.5, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
        _txt(s, l + Inches(0.48), y + Inches(0.06), w - Inches(0.60), box_h - Inches(0.12),
             text, sz=8.5, color=C_DARK)
    return y0 + len(items) * (box_h + gap)

def _arrow(s, l, t, size=Inches(0.3), color=C_LGREY):
    _txt(s, l, t, size, size, "▸", sz=16, color=color, bold=True, align=PP_ALIGN.CENTER)

def _label_row(s, l, t, label_w, val_w, label, value, row_h=Inches(0.42), bg=C_WHITE):
    _shape(s, l, t, label_w, row_h, C_PRIMARY)
    _txt(s, l + Inches(0.1), t + Inches(0.08), label_w - Inches(0.2), row_h - Inches(0.16),
         label, sz=9, color=C_WHITE, bold=True)
    _shape(s, l + label_w, t, val_w, row_h, bg, C_BORDER)
    _txt(s, l + label_w + Inches(0.15), t + Inches(0.08), val_w - Inches(0.3), row_h - Inches(0.16),
         value, sz=9, color=C_DARK)

def _insight_panel(s, l, t, w, h, title, bullets):
    actual_h = max(h, Inches(0.40 + len(bullets[:6]) * 0.22))
    _shape(s, l, t, w, actual_h, C_INSIGHT, C_BORDER, radius=True)
    _shape(s, l, t, Inches(0.06), actual_h, C_PRIMARY, radius=False)
    _txt(s, l + Inches(0.18), t + Inches(0.10), w - Inches(0.34), Inches(0.18),
         title, sz=9, color=C_PRIMARY, bold=True)
    for i, b in enumerate(bullets[:6]):
        _txt(s, l + Inches(0.18), t + Inches(0.35) + Inches(i * 0.22),
             w - Inches(0.34), Inches(0.20), f"· {b}", sz=9, color=C_DARK)

def _table(s, l, t, headers, rows, col_widths, row_h=Inches(0.32)):
    tbl = s.shapes.add_table(len(rows) + 1, len(headers), l, t,
                              sum(col_widths), row_h * (len(rows) + 1)).table
    for j, (hdr, cw) in enumerate(zip(headers, col_widths)):
        tbl.columns[j].width = cw
        cell = tbl.cell(0, j)
        cell.text = hdr
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(8); p.font.bold = True; p.font.color.rgb = C_WHITE
            p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
        cell.fill.solid(); cell.fill.fore_color.rgb = C_PRIMARY
        cell.margin_top = Pt(4); cell.margin_bottom = Pt(4)
        cell.margin_left = Pt(4); cell.margin_right = Pt(4)
    for i, row_data in enumerate(rows):
        bg = C_WHITE
        for j, val in enumerate(row_data):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8); p.font.color.rgb = C_BLACK
                p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            cell.margin_top = Pt(4); cell.margin_bottom = Pt(4)
            cell.margin_left = Pt(4); cell.margin_right = Pt(4)
    # 테이블 라인 스타일 — 얇은 회색 선
    try:
        from pptx.oxml.ns import qn
        tbl_elm = tbl._tbl
        tblPr = tbl_elm.find(qn('a:tblPr'))
        if tblPr is None:
            from lxml import etree
            tblPr = etree.SubElement(tbl_elm, qn('a:tblPr'))
        tblPr.set('bandRow', '0')
    except Exception:
        pass
    return tbl

def _bar_h(s, l, t, w, h, pct, fill=C_PRIMARY, bg=C_PALE):
    _shape(s, l, t, w, h, bg, radius=True)
    fw = max(int(w * min(pct, 1.0)), Inches(0.05))
    _shape(s, l, t, fw, h, fill, radius=True)

def _chart_img(s, fig, l, t, w_in, h_in):
    try:
        import kaleido  # noqa: F401 — 설치 확인
        tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        fig.write_image(tmp.name, width=int(w_in * 100), height=int(h_in * 100),
                        scale=2, engine="kaleido")
        s.shapes.add_picture(tmp.name, Inches(l), Inches(t), Inches(w_in))
        os.unlink(tmp.name)
        return True
    except Exception:
        # kaleido 미설치 또는 렌더링 실패 시 빈 자리에 안내 텍스트
        _shape(s, Inches(l), Inches(t), Inches(w_in), Inches(h_in), C_XPALE, C_BORDER, radius=True)
        _txt(s, Inches(l), Inches(t + h_in/2 - 0.15), Inches(w_in), Inches(0.30),
             "차트 렌더링 불가 (kaleido 패키지 필요)",
             sz=9, color=C_GREY, align=PP_ALIGN.CENTER)
        return False


# ══════════════════════════════════════════════════════
# 메인 생성 함수
# ══════════════════════════════════════════════════════
def generate_lp_pptx(
    summary: dict,
    result_df: pd.DataFrame,
    commentary: str,
    quarter: str = "",
    fund_name: str = "PE/VC 펀드",
    fund_strategy: str = "VC",
    base_date: str = "",
    selected_sections: list = None,
    include_charts: bool = False,
    jcurve_df=None,
    scenario_df=None, scenario_company: str = "",
    sensitivity_df=None, sensitivity_company: str = "",
    dart_fin_df=None, dart_company: str = "",
    kvic_sector_df=None,
    rate_df=None, fx_df=None, spread=None,
    df_raw=None,
) -> bytes:

    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    sel = set(selected_sections) if selected_sections else set()
    include_all = len(sel) == 0
    def _sec(kw):
        return include_all or any(kw in s for s in sel)

    # 자주 쓰는 값
    moic     = summary["펀드 MOIC"]
    dpi      = summary["펀드 DPI"]
    rvpi     = summary["펀드 RVPI"]
    tvpi     = summary["펀드 TVPI"]
    avg_irr  = round(result_df["IRR(%)"].mean(), 1)
    n_cos    = summary["포트폴리오사 수"]
    total_inv = summary["총 투자금액 (백만원)"]
    total_val = result_df["현재가치_백만원"].sum() + result_df["회수금액_백만원"].sum()

    slide_labels = []

    # ════════════════════════════════════════════════
    # 1. 표지 (항상 포함) — 상단 컬러블록 + 하단 흰색
    # ════════════════════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, C_WHITE)

    # 상단 초록 단색 블록 (슬라이드 약 88%) — 배경 이미지 사용하지 않음
    cover_h = Inches(6.6)
    _shape(s, Inches(0), Inches(0), SW, cover_h, C_PRIMARY)

    # 제목 (상단 블록 중앙)
    _txt(s, Inches(1.2), Inches(1.8), Inches(10), Inches(0.8),
         fund_name, sz=42, color=C_WHITE, bold=True)
    _txt(s, Inches(1.2), Inches(2.8), Inches(10), Inches(0.5),
         f"{fund_strategy} 포트폴리오 {quarter} 분기 성과 보고서", sz=18, color=C_PALE)

    # 하단 흰색 영역 — 정보
    _shape(s, Inches(0), cover_h, SW, SH - cover_h, C_WHITE)
    _shape(s, Inches(0), cover_h, SW, Pt(4), C_PRIMARY)

    _txt(s, Inches(1.2), cover_h + Inches(0.28), Inches(5), Inches(0.20),
         f"{base_date}  ·  이수빈  ·  SDIC", sz=11, color=C_DARK)
    _txt(s, Inches(1.2), cover_h + Inches(0.55), Inches(8), Inches(0.18),
         "본 보고서는 포트폴리오 성과 데이터를 기반으로 자동 생성되었습니다.", sz=8, color=C_LGREY)

    slide_labels.append("표지")

    # ════════════════════════════════════════════════
    # 2. Executive Summary
    # ════════════════════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
    _header(s, "Executive Summary", "Executive Summary",
            f"{fund_name}의 {quarter} 포트폴리오 성과를 종합 요약합니다.")

    # 핵심 지표 한줄 요약
    total_cur = result_df["현재가치_백만원"].sum()
    total_rec = result_df["회수금액_백만원"].sum()
    total_profit = (total_cur + total_rec) - total_inv
    profit_pct = total_profit / total_inv * 100 if total_inv > 0 else 0

    _txt(s, M_LEFT, BODY_Y, C_WIDTH, Inches(0.30),
         f"{n_cos}개 기업에 총 {total_inv:,}백만원을 투자하여 MOIC {moic}x, IRR {avg_irr}%의 성과를 달성하였습니다.",
         sz=13, color=C_BLACK, bold=True)

    # 레퍼런스 스타일: 좌측 라벨(다크그린) + 2개 값 컬럼 그리드
    bm_over = len(result_df[result_df["MOIC"] >= 2.0])
    bm_under = len(result_df[result_df["MOIC"] < 1.0])
    top_r = result_df.sort_values("MOIC", ascending=False).iloc[0]
    bot_r = result_df.sort_values("MOIC", ascending=True).iloc[0]
    realized_ratio = total_rec / (total_cur + total_rec) * 100 if (total_cur + total_rec) > 0 else 0

    es_y = BODY_Y + Inches(0.55)
    label_w = Inches(1.5)
    col_w = Inches(5.25)
    col2_x = M_LEFT + label_w + col_w

    # 헤더 행 (Core System / AI & Export 처럼 두 컬럼 제목)
    _shape(s, M_LEFT, es_y, label_w, Inches(0.32), C_WHITE)
    _shape(s, M_LEFT + label_w, es_y, col_w, Inches(0.32), C_BG, C_BORDER)
    _txt(s, M_LEFT + label_w, es_y + Inches(0.05), col_w, Inches(0.24), "성과 지표", sz=10, color=C_DARK, bold=True, align=PP_ALIGN.CENTER)
    _shape(s, col2_x, es_y, col_w, Inches(0.32), C_BG, C_BORDER)
    _txt(s, col2_x, es_y + Inches(0.05), col_w, Inches(0.24), "포트폴리오 현황", sz=10, color=C_DARK, bold=True, align=PP_ALIGN.CENTER)

    rows_y = es_y + Inches(0.32)

    def _es_row(y, h, label, left_lines, right_lines, highlight=False):
        bg = C_PALE if highlight else C_WHITE
        _shape(s, M_LEFT, y, label_w, h, C_PRIMARY)
        _txt(s, M_LEFT + Inches(0.1), y + Inches(0.10), label_w - Inches(0.2), h - Inches(0.2),
             label, sz=10, color=C_WHITE, bold=True)
        _shape(s, M_LEFT + label_w, y, col_w, h, bg, C_BORDER)
        _multiline(s, M_LEFT + label_w + Inches(0.12), y + Inches(0.08), col_w - Inches(0.24), h - Inches(0.16),
                   left_lines, sz=8.5, color=C_DARK, spacing=3)
        _shape(s, col2_x, y, col_w, h, bg, C_BORDER)
        _multiline(s, col2_x + Inches(0.12), y + Inches(0.08), col_w - Inches(0.24), h - Inches(0.16),
                   right_lines, sz=8.5, color=C_DARK, spacing=3)

    h1 = Inches(0.55)
    _es_row(rows_y, h1, "펀드 개요",
            [f"{fund_name} · {fund_strategy}", f"{quarter} · 기준일 {base_date}"],
            [f"포트폴리오 {n_cos}개 기업", f"총 투자금 {total_inv:,}백만원"])

    y2 = rows_y + h1
    h2 = Inches(0.55)
    _es_row(y2, h2, "핵심 성과",
            [f"MOIC {moic}x · IRR {avg_irr}%", f"BM 대비 {'달성' if moic >= 2.0 else f'{moic/2.0*100:.0f}% 수준'}"],
            [f"TVPI {tvpi}x · DPI {dpi}x", f"RVPI {rvpi}x"])

    y3 = y2 + h2
    h3 = Inches(0.85)
    _es_row(y3, h3, "가치 창출",
            [f"투자금 {total_inv:,.0f}M → 가치 {total_cur+total_rec:,.0f}M", f"전체 수익률 {profit_pct:+.0f}%", f"현금 실현율 {realized_ratio:.0f}%"],
            [f"BM 달성 {bm_over}개사 / 원금 미달 {bm_under}개사", f"최고: {top_r['회사명']}({top_r['MOIC']}x)", f"최저: {bot_r['회사명']}({bot_r['MOIC']}x)"],
            highlight=True)

    es_bot = y3 + h3 + Inches(0.20)

    # 하단: 기대 효과 4분할 (레퍼런스 "단기효과/중장기효과" 스타일)
    _txt(s, M_LEFT, es_bot, Inches(4), Inches(0.18), "향후 모니터링 포인트", sz=9, color=C_GREY, bold=True)
    box_w = Inches(2.85); box_h = Inches(1.1); box_gap = Inches(0.15)
    monitor_items = [
        ("성과 모니터링", f"MOIC {moic}x 유지/개선 여부를 분기별로 추적하며, BM(2.0x) 미달 시 원인 분석 필요"),
        ("리스크 관리", f"원금 미달 {bm_under}개사에 대한 집중 모니터링 및 Exit 전략 재검토"),
        ("회수 전략", f"실현율 {realized_ratio:.0f}% — {'후속 회수 가속화' if realized_ratio < 30 else '회수 모멘텀 유지'} 방안 수립"),
        ("포트폴리오 운영", f"섹터 분산 및 투자단계별 밸런스 점검을 통한 리스크 분산 강화"),
    ]
    for i, (mt, md) in enumerate(monitor_items):
        mx = M_LEFT + i * (box_w + box_gap)
        _process_box(s, mx, es_bot + Inches(0.22), box_w, box_h, i + 1, mt, md, bg=C_XPALE)

    slide_labels.append("Executive Summary")

    # ════════════════════════════════════════════════
    # 3. Agenda
    # ════════════════════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, C_WHITE)

    # 좌측: "Agenda" 타이틀 (세로 중앙)
    _txt(s, Inches(1.5), Inches(2.8), Inches(4), Inches(0.8),
         "Agenda", sz=40, color=C_DARK, bold=True)

    # 세로 구분선 (중앙) — 굵고 진한 검정
    _shape(s, Inches(6.3), Inches(1.2), Pt(3), Inches(5.2), C_BLACK)

    # 우측: 목차 항목 (세로 중앙 정렬, 우측으로 이동)
    agenda = ["펀드 성과 종합", "포트폴리오 · 섹터 분석", "Top/Bottom · 리스크",
              "J-Curve · 시나리오 · Sensitivity", "Waterfall 분배",
              "KVIC 시장 비교 · 거시지표", "AI 코멘터리", "Appendix"]
    agenda_keys = [["성과"], ["포트폴리오", "섹터"], ["Top", "리스크"],
                    ["J-Curve", "시나리오", "Sensitivity"], ["Waterfall"],
                    ["KVIC", "거시", "DART"], ["AI"], None]
    n_items = len(agenda)
    item_h = 0.55
    total_h = n_items * item_h
    start_y = (7.5 - total_h) / 2
    for i, item in enumerate(agenda):
        keys = agenda_keys[i] if i < len(agenda_keys) else [item.split()[0]]
        active = True if keys is None else any(any(k in sn for k in keys) for sn in (selected_sections or []))
        y = Inches(start_y + i * item_h)
        _txt(s, Inches(7.1), y, Inches(0.5), Inches(0.35),
             f"{i + 1}.", sz=14, color=C_DARK if active else C_LGREY,
             bold=active, align=PP_ALIGN.RIGHT)
        _txt(s, Inches(7.7), y, Inches(5), Inches(0.35),
             item, sz=14, color=C_DARK if active else C_LGREY, bold=active)
    slide_labels.append("Agenda")

    # ════════════════════════════════════════════════
    # 3. 성과 요약
    # ════════════════════════════════════════════════
    if _sec("성과 요약"):
        s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
        bm_msg = f"MOIC {moic}x · IRR {avg_irr}% — {'벤치마크(MOIC 2.0x, IRR 15%) 달성' if moic >= 2.0 and avg_irr >= 15 else '벤치마크 대비 성장 중'}"
        _header(s, "PERFORMANCE SUMMARY", "성과 요약", bm_msg)

        # ── ROW 1: 핵심 KPI 4개 ──
        r1 = BODY_Y + Inches(0.10)
        kpi_w = Inches(2.85); kpi_gap = Inches(0.15)
        for i, (lab, val, sub, vc) in enumerate([
            ("MOIC", f"{moic}x", "투자 대비 전체 가치", C_PRIMARY),
            ("IRR (평균)", f"{avg_irr}%", "연환산 수익률", C_PRIMARY),
            ("TVPI", f"{tvpi}x", "DPI + RVPI", C_PRIMARY if tvpi >= 2.0 else C_BLACK),
            ("투자기업 수", f"{n_cos}개", f"총 {total_inv:,}백만원", C_PRIMARY),
        ]):
            x = M_LEFT + (kpi_w + kpi_gap) * i
            _kpi_card(s, x, r1, kpi_w, Inches(0.95), lab, val, sub, vc)

        # 지표 정의 한 줄
        _txt(s, M_LEFT, r1 + Inches(1.00), C_WIDTH, Inches(0.16),
             "MOIC = 총가치/투자금  ·  IRR = 연환산수익률  ·  TVPI = DPI+RVPI  ·  DPI = 회수금/투자금  ·  RVPI = 미실현/투자금",
             sz=7, color=C_LGREY)

        # ── ROW 2: 가치 흐름 (투자→가치→회수→수익) ──
        r2 = r1 + Inches(1.25)
        _txt(s, M_LEFT, r2, Inches(5), Inches(0.18),
             "VALUE CREATION FLOW", sz=9, color=C_GREY, bold=True)

        total_cur = result_df["현재가치_백만원"].sum()
        total_rec = result_df["회수금액_백만원"].sum()
        total_profit = (total_cur + total_rec) - total_inv
        profit_pct = total_profit / total_inv * 100 if total_inv > 0 else 0

        flow_items = [
            ("투자금", f"{total_inv:,.0f}M", f"{n_cos}개사에 집행", C_DARK),
            ("현재가치", f"{total_cur:,.0f}M", f"미실현 NAV", C_ACCENT),
            ("회수금", f"{total_rec:,.0f}M", f"현금 실현", C_PRIMARY),
            ("수익", f"{total_profit:+,.0f}M", f"수익률 {profit_pct:+.0f}%", C_PRIMARY if total_profit >= 0 else C_RED),
        ]
        fw = Inches(2.65)
        arrow_w = Inches(0.35)
        # 컬러별 배경 매핑: 투자금(베이지) / 현재가치(연초록) / 회수금(초록) / 수익(초록 or 연빨)
        _flow_bg = {
            C_DARK: C_BEIGE, C_ACCENT: C_XPALE, C_PRIMARY: C_PALE,
            C_RED: RGBColor(0xFF, 0xEB, 0xEE),
        }
        for i, (fl, fv, fd, fc) in enumerate(flow_items):
            fx = M_LEFT + i * (fw + arrow_w)
            bg_c = _flow_bg.get(fc, C_WHITE)
            _shape(s, fx, r2 + Inches(0.25), fw, Inches(0.70), bg_c, C_BORDER, radius=True)
            _shape(s, fx, r2 + Inches(0.25), fw, Pt(3.5), fc)   # 상단 컬러 스트라이프
            _txt(s, fx, r2 + Inches(0.32), fw, Inches(0.15), fl, sz=8, color=C_GREY, bold=True, align=PP_ALIGN.CENTER)
            _txt(s, fx, r2 + Inches(0.49), fw, Inches(0.25), fv, sz=16, color=fc, bold=True, align=PP_ALIGN.CENTER)
            _txt(s, fx, r2 + Inches(0.72), fw, Inches(0.15), fd, sz=7, color=C_GREY, align=PP_ALIGN.CENTER)
            if i < 3:
                ax = fx + fw + Inches(0.05)
                _txt(s, ax, r2 + Inches(0.40), Inches(0.25), Inches(0.28), "▶", sz=14, color=fc, bold=True, align=PP_ALIGN.CENTER)

        # ── ROW 3: 좌측(DPI/RVPI + 벤치마크) / 우측(등급분포 + 인사이트) ──
        r3 = r2 + Inches(1.45)   # 위 섹션과 충분한 간격
        half = Inches(5.85)

        # 좌측: DPI/RVPI 분해 + 벤치마크 바
        _kpi_card(s, M_LEFT, r3, Inches(2.8), Inches(0.75), "DPI (현금회수)", f"{dpi}x", "회수금 / 투자금", C_PRIMARY)
        _kpi_card(s, M_LEFT + Inches(2.95), r3, Inches(2.8), Inches(0.75), "RVPI (잔존가치)", f"{rvpi}x", "미실현 / 투자금", C_ACCENT)

        bm_y = r3 + Inches(0.95)
        _txt(s, M_LEFT, bm_y, half, Inches(0.18), "벤치마크 달성률", sz=9, color=C_GREY, bold=True)
        bm_criteria = ["우수 기준 2.0x (업계 중위수)", "목표 수익률 15% (VC 평균)", "DPI+RVPI 합산 2.0x"]
        for j, (nm, act, tgt, crit) in enumerate([
            ("MOIC", float(moic), 2.0, bm_criteria[0]),
            ("IRR",  float(avg_irr), 15.0, bm_criteria[1]),
            ("TVPI", float(tvpi), 2.0, bm_criteria[2]),
        ]):
            by = bm_y + Inches(0.22) + Inches(j * 0.30)
            _txt(s, M_LEFT, by, Inches(0.6), Inches(0.20), nm, sz=9, color=C_BLACK, bold=True)
            pct = min(act / tgt, 1.5) / 1.5
            _bar_h(s, M_LEFT + Inches(0.65), by + Inches(0.04), Inches(3.4), Inches(0.14), pct)
            clr = C_PRIMARY if act >= tgt else C_RED
            _txt(s, M_LEFT + Inches(4.10), by, Inches(0.65), Inches(0.20), f"{act}", sz=9, color=clr, bold=True)
            _txt(s, M_LEFT + Inches(4.75), by, Inches(0.7), Inches(0.20), f"/ {tgt}", sz=8, color=C_GREY)
            _txt(s, M_LEFT + Inches(0.65), by + Inches(0.19), Inches(3.4), Inches(0.12),
                 crit, sz=6, color=C_LGREY)

        # 우측: 포트폴리오 등급 분포
        rx = M_LEFT + Inches(6.3)
        rw = Inches(5.8)
        _txt(s, rx, r3, rw, Inches(0.18), "포트폴리오 등급 분포", sz=9, color=C_GREY, bold=True)

        grade_counts = {"S (3.0x↑)": 0, "A (2.0x)": 0, "B (1.5x)": 0, "C (1.0x)": 0, "D (<1.0x)": 0}
        for m in result_df["MOIC"]:
            if m >= 3.0: grade_counts["S (3.0x↑)"] += 1
            elif m >= 2.0: grade_counts["A (2.0x)"] += 1
            elif m >= 1.5: grade_counts["B (1.5x)"] += 1
            elif m >= 1.0: grade_counts["C (1.0x)"] += 1
            else: grade_counts["D (<1.0x)"] += 1

        grade_criteria = ["MOIC 3.0x+", "2.0~3.0x", "1.5~2.0x", "1.0~1.5x", "<1.0x"]
        gcolors = [C_PRIMARY, C_SECONDARY, C_ACCENT, C_ORANGE, C_RED]
        _g_start = r3 + Inches(0.36)   # 라벨과 카드 사이 여유 확보
        for gi, ((gname, gcnt), gclr, gcrit) in enumerate(zip(grade_counts.items(), gcolors, grade_criteria)):
            gx = rx + Inches(gi * 1.15)
            _shape(s, gx, _g_start, Inches(1.05), Inches(0.72), C_WHITE, C_BORDER, radius=True)
            _shape(s, gx, _g_start, Inches(1.05), Pt(3.5), gclr)
            _txt(s, gx, _g_start + Inches(0.08), Inches(1.05), Inches(0.14), gname,
                 sz=7.5, color=gclr, bold=True, align=PP_ALIGN.CENTER)
            _txt(s, gx, _g_start + Inches(0.23), Inches(1.05), Inches(0.13), gcrit,
                 sz=6.5, color=C_GREY, align=PP_ALIGN.CENTER)
            _txt(s, gx, _g_start + Inches(0.37), Inches(1.05), Inches(0.22), f"{gcnt}개사",
                 sz=13, color=gclr, bold=True, align=PP_ALIGN.CENTER)
            # 비중 미니 바
            _bar_h(s, gx + Inches(0.10), _g_start + Inches(0.60), Inches(0.85), Inches(0.06),
                   gcnt / n_cos if n_cos > 0 else 0, fill=gclr, bg=C_PALE)

        # 우측 하단: KEY INSIGHT
        moic_over2 = len(result_df[result_df["MOIC"] >= 2.0])
        moic_under1 = len(result_df[result_df["MOIC"] < 1.0])
        realized_ratio = total_rec / (total_cur + total_rec) * 100 if (total_cur + total_rec) > 0 else 0

        insights = []
        if moic >= 2.0:
            insights.append(f"MOIC {moic}x로 벤치마크(2.0x) 달성 — 상위 펀드 수준")
        else:
            insights.append(f"MOIC {moic}x — 벤치마크 대비 {moic/2.0*100:.0f}% 달성, {2.0-moic:.2f}x 추가 성장 필요")
        insights.append(f"투자 {total_inv:,.0f}M → 가치 {total_cur+total_rec:,.0f}M 창출 (수익률 {profit_pct:+.0f}%)")
        insights.append(f"현금 실현율 {realized_ratio:.0f}% — {'회수 진행 중' if realized_ratio > 30 else '초기 투자 단계, 미실현 가치 위주'}")
        insights.append(f"BM 달성 {moic_over2}개사 / 원금 미달 {moic_under1}개사 — 전체 {n_cos}개 포트폴리오")
        _insight_panel(s, rx, _g_start + Inches(0.82), rw, Inches(1.10), "성과 종합 코멘트", insights)

        slide_labels.append("성과")

    # ════════════════════════════════════════════════
    # 4. 포트폴리오 상세
    # ════════════════════════════════════════════════
    if _sec("포트폴리오"):
        s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
        top1 = result_df.sort_values("MOIC", ascending=False).iloc[0]
        bot1 = result_df.sort_values("MOIC", ascending=True).iloc[0]
        _header(s, "PORTFOLIO DETAIL", "포트폴리오 상세",
                f'{n_cos}개 포트폴리오 기업 중 {top1["회사명"]}이 MOIC {top1["MOIC"]}x로 가장 높은 성과를 보이고 있습니다.')

        # 전체 테이블 (최대 10개사)
        sorted_df = result_df.sort_values("MOIC", ascending=False).head(10)
        hdrs = ["회사명", "섹터", "투자단계", "투자(M)", "현재가치(M)", "회수(M)", "MOIC", "IRR(%)", "DPI", "RVPI", "TVPI"]
        cw = [Inches(1.3), Inches(0.9), Inches(0.75), Inches(0.85), Inches(0.95), Inches(0.85),
              Inches(0.65), Inches(0.7), Inches(0.55), Inches(0.55), Inches(0.55)]
        rows = []
        for _, r in sorted_df.iterrows():
            rows.append([
                r["회사명"], r.get("섹터", "-"), r.get("투자단계", "-"),
                f'{int(r["투자금액_백만원"]):,}', f'{int(r["현재가치_백만원"]):,}',
                f'{int(r["회수금액_백만원"]):,}',
                f'{r["MOIC"]}x', f'{r["IRR(%)"]}%',
                f'{r["DPI"]}x', f'{r["RVPI"]}x', f'{r["TVPI"]}x',
            ])
        _table(s, M_LEFT, BODY_Y, hdrs, rows, cw, Inches(0.28))

        # 차트 (테이블 아래 좌측)
        chart_y = BODY_Y + Inches(0.28) * (len(rows) + 1) + Inches(0.20)
        moic_over2 = len(result_df[result_df["MOIC"] >= 2.0])
        moic_under1 = len(result_df[result_df["MOIC"] < 1.0])
        avg_moic = result_df["MOIC"].mean()
        med_moic = result_df["MOIC"].median()
        top_name = result_df.sort_values("MOIC", ascending=False).iloc[0]["회사명"]
        top_moic = result_df.sort_values("MOIC", ascending=False).iloc[0]["MOIC"]
        top_sector_name = result_df.sort_values("MOIC", ascending=False).iloc[0].get("섹터", "-")

        if include_charts:
            import plotly.graph_objects as go
            sd = result_df.sort_values("MOIC", ascending=True)
            fig = go.Figure(go.Bar(
                x=sd["MOIC"].tolist(), y=sd["회사명"].tolist(), orientation="h",
                marker_color=["rgba(27,94,32,0.7)" if m >= 2 else "rgba(67,160,71,0.5)" if m >= 1 else "rgba(198,40,40,0.4)" for m in sd["MOIC"]],
                text=[f"{m}x" for m in sd["MOIC"]], textposition="outside", textfont=dict(size=8),
                marker_line_width=0,
            ))
            fig.add_vline(x=2.0, line_dash="dot", line_color="#999", annotation_text="BM 2.0x", annotation_font_size=8)
            fig.update_layout(height=270, width=560, margin=dict(t=10, b=45, l=80, r=40),
                              paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)",
                              xaxis=dict(showgrid=True, gridcolor="#eee", title="MOIC (배)"), yaxis=dict(showgrid=False), bargap=0.3)
            _chart_img(s, fig, 0.6, min(float(chart_y / 914400), 4.6), 6.3, 2.95)

        # 우측: 줄글 해설 (차트 옆 빈 공간 채움) — 헤더 바 포함 카드형
        narrative = (
            f"본 펀드는 {n_cos}개 포트폴리오 기업에 총 {total_inv:,}백만원을 투자하였으며, "
            f"이 중 {moic_over2}개사({moic_over2/n_cos*100:.0f}%)가 벤치마크 MOIC 2.0x를 상회하는 성과를 기록하고 있습니다. "
            f"가장 우수한 성과를 보인 {top_name}({top_sector_name})은 MOIC {top_moic}x를 달성하며 펀드 전체 수익의 핵심 동력으로 작용하고 있습니다.\n\n"
            f"전체 포트폴리오의 평균 MOIC는 {avg_moic:.2f}x, 중앙값은 {med_moic:.2f}x로, "
            f"{'평균이 중앙값을 크게 상회하여 소수 우수 기업에 수익이 집중되는 구조' if avg_moic > med_moic * 1.2 else '평균과 중앙값의 차이가 크지 않아 비교적 고른 성과 분포'}를 보이고 있습니다. "
            f"{f'한편 원금을 회수하지 못한 기업이 {moic_under1}개사 존재하여, 해당 기업에 대한 후속 모니터링과 Exit 전략 재검토가 필요합니다.' if moic_under1 > 0 else '모든 포트폴리오 기업이 투자 원금 이상의 가치를 유지하고 있어 안정적인 펀드 운용 상태입니다.'}"
        )
        ins_y = min(float(chart_y / 914400), 4.6)
        narr_w = Inches(5.6)
        narr_h = Inches(2.15)
        _shape(s, Inches(7.1), Inches(ins_y), narr_w, narr_h, C_WHITE, C_BORDER, radius=True)
        _shape(s, Inches(7.1), Inches(ins_y), narr_w, Inches(0.34), C_PRIMARY, radius=True)
        _txt(s, Inches(7.1), Inches(ins_y) + Inches(0.07), narr_w, Inches(0.22),
             "PORTFOLIO NARRATIVE", sz=9, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
        _txt(s, Inches(7.3), Inches(ins_y) + Inches(0.48), narr_w - Inches(0.4), narr_h - Inches(0.60), narrative, sz=9, color=C_DARK)
        slide_labels.append("포트폴리오")

    # ════════════════════════════════════════════════
    # 5. 섹터 분석
    # ════════════════════════════════════════════════
    if _sec("섹터") and "섹터" in result_df.columns:
        s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
        sa = result_df.groupby("섹터").agg(
            기업수=("회사명", "count"), 총투자=("투자금액_백만원", "sum"),
            평균MOIC=("MOIC", "mean"), 평균IRR=("IRR(%)", "mean")
        ).sort_values("총투자", ascending=False).reset_index()
        total_all = sa["총투자"].sum()
        top_sec = sa.iloc[0]
        best_moic_sec = sa.sort_values("평균MOIC", ascending=False).iloc[0]
        top3_pct = sa.head(3)["총투자"].sum() / total_all * 100
        _header(s, "SECTOR ANALYSIS", "섹터별 투자 분석",
                f"총 {len(sa)}개 섹터에 분산 투자하며, {top_sec['섹터']}에 전체의 {top_sec['총투자']/total_all*100:.0f}%가 집중되어 있습니다.")

        # ── 페이지 1: 테이블 + 차트 + KPI + 줄글 ──
        sec_hdrs = ["섹터", "기업수", "투자(M)", "비중", "MOIC", "IRR"]
        sec_cw = [Inches(1.1), Inches(0.55), Inches(0.9), Inches(0.6), Inches(0.7), Inches(0.7)]
        sec_rows = []
        for _, r in sa.head(8).iterrows():
            pct = r["총투자"] / total_all * 100 if total_all > 0 else 0
            sec_rows.append([r["섹터"], f'{int(r["기업수"])}', f'{int(r["총투자"]):,}',
                             f'{pct:.0f}%', f'{r["평균MOIC"]:.1f}x', f'{r["평균IRR"]:.0f}%'])
        _table(s, M_LEFT, BODY_Y, sec_hdrs, sec_rows, sec_cw, Inches(0.32))

        # 우측: Plotly 바 차트 — 테이블과 간격을 넓혀 배치
        chart_x = 6.5
        if include_charts:
            import plotly.graph_objects as go
            sa_plot = sa.head(6).sort_values("총투자", ascending=True)
            fig = go.Figure(go.Bar(
                x=sa_plot["총투자"].tolist(), y=sa_plot["섹터"].tolist(), orientation="h",
                marker_color=["#1b5e20", "#2e7d32", "#43a047", "#66bb6a", "#a5d6a7", "#c8e6c9"][:len(sa_plot)],
                text=[f'{int(v):,}M ({v/total_all*100:.0f}%)' for v in sa_plot["총투자"]],
                textposition="outside", textfont=dict(size=8), marker_line_width=0,
            ))
            fig.update_layout(height=250, width=460, margin=dict(t=5, b=30, l=80, r=70),
                              paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)",
                              xaxis=dict(showgrid=True, gridcolor="#eee", title="투자금(백만원)"),
                              yaxis=dict(showgrid=False), bargap=0.3)
            _chart_img(s, fig, chart_x, float(BODY_Y / 914400) + 0.2, 5.4, 3.05)

        # 하단: KPI 카드 4개 (섹터 다변화 지표)
        kpi_y = BODY_Y + Inches(3.6)
        kw4 = Inches(2.85)
        gini_like = (sa["총투자"] / total_all).pow(2).sum()
        _kpi_card(s, M_LEFT, kpi_y, kw4, Inches(0.85), "섹터 수", f"{len(sa)}개", "분산 투자 범위", C_PRIMARY)
        _kpi_card(s, M_LEFT + Inches(3.05), kpi_y, kw4, Inches(0.85), "상위 3개 집중도", f"{top3_pct:.0f}%", "Top3 섹터 비중", C_RED if top3_pct > 70 else C_PRIMARY)
        _kpi_card(s, M_LEFT + Inches(6.10), kpi_y, kw4, Inches(0.85), "최고 MOIC 섹터", f"{best_moic_sec['평균MOIC']:.1f}x", best_moic_sec["섹터"], C_PRIMARY)
        _kpi_card(s, M_LEFT + Inches(9.15), kpi_y, kw4, Inches(0.85), "섹터 집중지수", f"{gini_like:.2f}", "HHI식 분산도(0~1)", C_ORANGE if gini_like > 0.3 else C_PRIMARY)

        # 하단: 줄글 분석
        sec_narrative = (
            f"본 펀드는 총 {len(sa)}개 섹터에 {n_cos}개 포트폴리오 기업을 분산 배치하였으며, "
            f"이 중 {top_sec['섹터']} 섹터에 {top_sec['총투자']:,.0f}백만원({top_sec['총투자']/total_all*100:.0f}%)이 집중되어 가장 큰 투자 비중을 차지하고 있습니다. "
            f"상위 3개 섹터({', '.join(sa.head(3)['섹터'].tolist())})가 전체 투자금의 {top3_pct:.0f}%를 차지하여, "
            f"{'섹터 편중도가 다소 높은 편이며 향후 신규 투자 시 미투자 섹터로의 분산을 검토할 필요가 있습니다.' if top3_pct > 70 else '비교적 고른 섹터 분산을 유지하고 있습니다.'} "
            f"수익성 측면에서는 {best_moic_sec['섹터']} 섹터가 평균 MOIC {best_moic_sec['평균MOIC']:.1f}x, IRR {best_moic_sec['평균IRR']:.0f}%로 가장 우수한 성과를 보이고 있어, "
            f"향후 동일 섹터 내 후속 투자 기회를 적극적으로 검토할 가치가 있습니다."
        )
        sec_sum_y = kpi_y + Inches(1.05)
        _shape(s, M_LEFT, sec_sum_y, C_WIDTH, Inches(0.78), C_BG, C_BORDER, radius=True)
        _txt(s, M_LEFT + Inches(0.15), sec_sum_y + Inches(0.08), C_WIDTH - Inches(0.3), Inches(0.62),
             sec_narrative, sz=9, color=C_DARK)
        slide_labels.append("섹터")

        # ── 페이지 2: 섹터 심층 분석 ──
        s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
        _header(s, "SECTOR DEEP DIVE", "섹터 심층 분석",
                f"{best_moic_sec['섹터']} 섹터가 평균 MOIC {best_moic_sec['평균MOIC']:.1f}x로 가장 높은 수익률을 보이고 있습니다.")

        # 좌측: 분석 과정
        _txt(s, M_LEFT, BODY_Y, Inches(5), Inches(0.20), "SECTOR ANALYSIS PROCESS", sz=10, color=C_GREY, bold=True)
        process_items = [
            f"① 투자 분산 현황: {len(sa)}개 섹터에 {n_cos}개 기업 배분",
            f"② 집중도 분석: 상위 3개 섹터({', '.join(sa.head(3)['섹터'].tolist())})가 전체의 {top3_pct:.0f}%를 차지",
            f"③ 성과 비교: {best_moic_sec['섹터']} 섹터가 평균 MOIC {best_moic_sec['평균MOIC']:.1f}x, IRR {best_moic_sec['평균IRR']:.0f}%로 최고 성과",
            f"④ 최대 투자 섹터: {top_sec['섹터']}에 {top_sec['총투자']:,.0f}M ({top_sec['총투자']/total_all*100:.0f}%) 집중",
            f"⑤ 분산 평가: {'분산 양호 — 개별 섹터 부진 시에도 포트폴리오 전체 영향 제한적' if len(sa) >= 4 else '분산 부족 — 추가 섹터 다변화를 통한 리스크 분산 필요'}",
        ]
        _multiline(s, M_LEFT, BODY_Y + Inches(0.30), Inches(5.8), Inches(1.8), process_items, sz=10, color=C_DARK, spacing=7)

        # 우측: 섹터별 성과 비교 KPI
        rx = M_LEFT + Inches(6.3)
        _txt(s, rx, BODY_Y, Inches(5), Inches(0.20), "섹터별 성과 비교", sz=10, color=C_GREY, bold=True)
        for i, (_, r) in enumerate(sa.head(4).iterrows()):
            cy = BODY_Y + Inches(0.30) + Inches(i * 0.45)
            pct = r["총투자"] / total_all if total_all > 0 else 0
            _shape(s, rx, cy, Inches(5.8), Inches(0.38), C_WHITE, C_BORDER, radius=True)
            _txt(s, rx + Inches(0.1), cy + Inches(0.07), Inches(1.2), Inches(0.24),
                 r["섹터"], sz=10, color=C_BLACK, bold=True)
            _bar_h(s, rx + Inches(1.4), cy + Inches(0.10), Inches(2.0), Inches(0.15), pct)
            _txt(s, rx + Inches(3.5), cy + Inches(0.07), Inches(0.7), Inches(0.24),
                 f'{pct*100:.0f}%', sz=9, color=C_PRIMARY, bold=True)
            _txt(s, rx + Inches(4.2), cy + Inches(0.07), Inches(1.5), Inches(0.24),
                 f'MOIC {r["평균MOIC"]:.1f}x · IRR {r["평균IRR"]:.0f}%', sz=8, color=C_DARK)

        # 하단: 인사이트
        sec_ins = [
            f"최대 투자 섹터 {top_sec['섹터']}에 {top_sec['총투자']/total_all*100:.0f}% 집중 — {'과도한 편중으로 분산 필요' if top_sec['총투자']/total_all > 0.4 else '적정 수준의 집중도'}",
            f"최고 수익 섹터 {best_moic_sec['섹터']}는 MOIC {best_moic_sec['평균MOIC']:.1f}x, IRR {best_moic_sec['평균IRR']:.0f}%로 펀드 수익을 견인합니다.",
            f"{'투자 대비 성과가 낮은 섹터에 대한 후속 투자 검토 및 Exit 전략 수립이 필요합니다.' if sa.iloc[-1]['평균MOIC'] < 1.5 else '전반적으로 섹터별 성과가 양호한 수준입니다.'}",
        ]
        _insight_panel(s, M_LEFT, BODY_Y + Inches(2.5), C_WIDTH, Inches(1.3), "SECTOR INSIGHT", sec_ins)
        slide_labels.append("섹터 심층")

    # ════════════════════════════════════════════════
    # 6. Top / Bottom
    # ════════════════════════════════════════════════
    if _sec("Top/Bottom"):
        s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
        sd = result_df.sort_values("MOIC", ascending=False)
        _top_gap = sd.iloc[0]["MOIC"] - sd.iloc[-1]["MOIC"]
        _header(s, "TOP / BOTTOM", "성과 상위 · 하위 분석",
                f'{sd.iloc[0]["회사명"]}이 MOIC {sd.iloc[0]["MOIC"]}x로 최고 성과를 기록하였으며, 상하위 간 {_top_gap:.1f}x의 격차가 존재합니다.')

        half_w = Inches(5.85)
        # 좌측: Top 3
        _shape(s, M_LEFT, BODY_Y, half_w, Inches(0.28), C_PRIMARY)
        _txt(s, M_LEFT + Inches(0.1), BODY_Y + Inches(0.04), Inches(3), Inches(0.2),
             "▲ TOP PERFORMERS", sz=9, color=C_WHITE, bold=True)
        for i, (_, r) in enumerate(sd.head(3).iterrows()):
            y = BODY_Y + Inches(0.35) + Inches(i * 0.75)
            _shape(s, M_LEFT, y, half_w, Inches(0.65), C_WHITE, C_BORDER, radius=True)
            _shape(s, M_LEFT + Inches(0.1), y + Inches(0.12), Inches(0.35), Inches(0.35), C_PRIMARY, radius=True)
            _txt(s, M_LEFT + Inches(0.1), y + Inches(0.15), Inches(0.35), Inches(0.3),
                 str(i + 1), sz=12, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
            _txt(s, M_LEFT + Inches(0.55), y + Inches(0.08), Inches(2.0), Inches(0.2),
                 r["회사명"], sz=13, color=C_BLACK, bold=True)
            _txt(s, M_LEFT + Inches(0.55), y + Inches(0.32), Inches(2.0), Inches(0.18),
                 f'{r.get("섹터", "")} · {r.get("투자단계", "")}', sz=8, color=C_GREY)
            _txt(s, M_LEFT + Inches(3.0), y + Inches(0.12), Inches(1.0), Inches(0.2),
                 f'MOIC {r["MOIC"]}x', sz=12, color=C_PRIMARY, bold=True)
            _txt(s, M_LEFT + Inches(4.1), y + Inches(0.12), Inches(0.8), Inches(0.2),
                 f'IRR {r["IRR(%)"]}%', sz=10, color=C_DARK)
            _txt(s, M_LEFT + Inches(4.9), y + Inches(0.12), Inches(0.8), Inches(0.2),
                 f'TVPI {r["TVPI"]}x', sz=10, color=C_DARK)

        # 우측: Bottom 3
        rx = M_LEFT + half_w + Inches(0.43)
        _shape(s, rx, BODY_Y, half_w, Inches(0.28), C_RED)
        _txt(s, rx + Inches(0.1), BODY_Y + Inches(0.04), Inches(3), Inches(0.2),
             "▼ UNDERPERFORMERS", sz=9, color=C_WHITE, bold=True)
        for i, (_, r) in enumerate(sd.tail(3).iterrows()):
            y = BODY_Y + Inches(0.35) + Inches(i * 0.75)
            clr = C_RED if r["MOIC"] < 1.0 else C_GREY
            _shape(s, rx, y, half_w, Inches(0.65), C_WHITE, C_BORDER, radius=True)
            _shape(s, rx + Inches(0.1), y + Inches(0.12), Inches(0.35), Inches(0.35), clr, radius=True)
            _txt(s, rx + Inches(0.1), y + Inches(0.15), Inches(0.35), Inches(0.3),
                 str(i + 1), sz=12, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
            _txt(s, rx + Inches(0.55), y + Inches(0.08), Inches(2.0), Inches(0.2),
                 r["회사명"], sz=13, color=C_BLACK, bold=True)
            _txt(s, rx + Inches(0.55), y + Inches(0.32), Inches(2.0), Inches(0.18),
                 f'{r.get("섹터", "")} · {r.get("투자단계", "")}', sz=8, color=C_GREY)
            _txt(s, rx + Inches(3.0), y + Inches(0.12), Inches(1.0), Inches(0.2),
                 f'MOIC {r["MOIC"]}x', sz=12, color=clr, bold=True)
            _txt(s, rx + Inches(4.1), y + Inches(0.12), Inches(0.8), Inches(0.2),
                 f'IRR {r["IRR(%)"]}%', sz=10, color=C_DARK)
            _txt(s, rx + Inches(4.9), y + Inches(0.12), Inches(0.8), Inches(0.2),
                 f'TVPI {r["TVPI"]}x', sz=10, color=C_DARK)

        # 하단: 분석 과정(좌, 번호박스) + 인사이트(우)
        tb_bot = BODY_Y + Inches(3.1)
        tb_process = [
            f"MOIC 기준 전체 {n_cos}개사를 정렬하여 상위 3개와 하위 3개를 추출",
            f"상위 3개 평균 {sd.head(3)['MOIC'].mean():.1f}x vs 하위 3개 {sd.tail(3)['MOIC'].mean():.1f}x — 격차 {sd.head(3)['MOIC'].mean()-sd.tail(3)['MOIC'].mean():.1f}x",
            f"원금 미달(MOIC<1.0x) {len(sd[sd['MOIC']<1.0])}개사 → {'리스크 관리 필요' if len(sd[sd['MOIC']<1.0]) > 0 else '전 기업 원금 이상'}",
        ]
        _process_list(s, M_LEFT, tb_bot, half_w, tb_process, box_h=Inches(0.42), gap=Inches(0.08))

        tb_ins = [
            f"상위 기업이 펀드 수익을 견인하고 있으며, 하위 기업의 Exit 전략 재검토가 필요합니다.",
            f"{'원금 미달 기업에 대한 집중 모니터링이 시급합니다.' if len(sd[sd['MOIC']<1.0]) > 0 else '전 기업이 원금 이상을 유지하고 있어 안정적입니다.'}",
            f"성과 편차(표준편차 {result_df['MOIC'].std():.2f}) — {'분산이 커 상위 기업 의존도가 높습니다.' if result_df['MOIC'].std() > 1.0 else '안정적인 분포를 보이고 있습니다.'}",
        ]
        _insight_panel(s, rx, tb_bot + Inches(0.24), half_w, Inches(1.1), "PERFORMANCE GAP", tb_ins)
        slide_labels.append("Top/Bottom")

    # ════════════════════════════════════════════════
    # 7. 리스크 · 집중도
    # ════════════════════════════════════════════════
    if _sec("리스크") or _sec("집중도"):
        s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
        weights = result_df["투자금액_백만원"] / result_df["투자금액_백만원"].sum()
        hhi = round((weights ** 2).sum() * 10000)
        hhi_label = "HIGH" if hhi > 2500 else ("MEDIUM" if hhi > 1500 else "LOW")
        n_under = len(result_df[result_df["MOIC"] < 1.0])
        _header(s, "RISK & CONCENTRATION", "리스크 · 집중도 평가",
                f'포트폴리오 집중도(HHI)는 {hhi:,}으로 {hhi_label} 수준이며, 원금 미달 기업은 {n_under}개사입니다.')

        # 좌측: 리스크 항목
        risks = []
        under = result_df[result_df["MOIC"] < 1.0]
        if len(under) > 0:
            risks.append(("원금 손실 위험", f"MOIC<1.0x: {', '.join(under['회사명'].tolist())}", "HIGH"))
        if hhi > 2500:
            risks.append(("집중 리스크", f"HHI {hhi:,} — 특정 기업에 투자 편중", "HIGH"))
        elif hhi > 1500:
            risks.append(("집중도 보통", f"HHI {hhi:,} — 분산 확대 검토 필요", "MEDIUM"))
        else:
            risks.append(("분산 양호", f"HHI {hhi:,} — 적정 분산 수준", "LOW"))
        if dpi < 0.5:
            risks.append(("회수 지연", f"DPI {dpi}x — 현금 회수 제한적", "MEDIUM"))
        if moic >= 2.0:
            risks.append(("우수 성과", f"MOIC {moic}x — 벤치마크 달성", "POSITIVE"))

        # ── 좌측: 리스크 항목 (간격 줄임) ──
        left_w = Inches(5.85)
        _txt(s, M_LEFT, BODY_Y, Inches(4), Inches(0.18), "1. 리스크 항목 평가", sz=9, color=C_DARK, bold=True)
        for i, (title, desc, level) in enumerate(risks):
            y = BODY_Y + Inches(0.34) + Inches(i * 0.52)
            if level == "HIGH":
                icon_c, bg_c, bd_c = C_RED, RGBColor(0xFE, 0xF5, 0xF5), RGBColor(0xF0, 0xD0, 0xD0)
            elif level == "MEDIUM":
                icon_c, bg_c, bd_c = C_ORANGE, RGBColor(0xFF, 0xF8, 0xEE), RGBColor(0xF0, 0xE0, 0xC0)
            else:
                icon_c, bg_c, bd_c = C_PRIMARY, C_XPALE, C_PALE
            _shape(s, M_LEFT, y, left_w, Inches(0.44), bg_c, bd_c, radius=True)
            _shape(s, M_LEFT + Inches(0.10), y + Inches(0.07), Inches(0.28), Inches(0.28), icon_c, radius=True)
            icon_txt = "!" if level in ("HIGH", "MEDIUM") else "✓"
            _txt(s, M_LEFT + Inches(0.10), y + Inches(0.08), Inches(0.28), Inches(0.26),
                 icon_txt, sz=9, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
            _txt(s, M_LEFT + Inches(0.50), y + Inches(0.05), Inches(2.5), Inches(0.18),
                 title, sz=10, color=C_BLACK, bold=True)
            _txt(s, M_LEFT + Inches(0.50), y + Inches(0.24), left_w - Inches(0.65), Inches(0.16),
                 desc, sz=8, color=C_DARK)

        # ── 우측: KPI 3개 (한 행) ──
        rx = M_LEFT + Inches(6.3)
        realized_pct = round(result_df["회수금액_백만원"].sum() / total_val * 100, 1) if total_val > 0 else 0
        avg_days = 0
        if "투자일" in result_df.columns:
            try:
                avg_days = (pd.to_datetime(result_df.iloc[0].get("기준일", datetime.date.today())) - pd.to_datetime(result_df["투자일"].min())).days
            except Exception:
                pass
        avg_yrs = round(avg_days / 365.25, 1) if avg_days > 0 else 0

        _txt(s, rx, BODY_Y, Inches(4), Inches(0.18), "2. 핵심 지표", sz=9, color=C_GREY, bold=True)
        kpi_y2 = BODY_Y + Inches(0.22)
        _kpi_card(s, rx, kpi_y2, Inches(1.85), Inches(0.80), "HHI", f"{hhi:,}", hhi_label,
                  C_RED if hhi > 2500 else C_ORANGE if hhi > 1500 else C_PRIMARY)
        _kpi_card(s, rx + Inches(1.95), kpi_y2, Inches(1.85), Inches(0.80), "실현율", f"{realized_pct}%", "회수/전체가치",
                  C_PRIMARY if realized_pct > 50 else C_BLACK)
        _kpi_card(s, rx + Inches(3.90), kpi_y2, Inches(1.85), Inches(0.80), "보유기간", f"{avg_yrs}년", "최초투자~기준일", C_BLACK)

        # 우측: 투자단계 분포 (KPI 바로 아래)
        if "투자단계" in result_df.columns:
            _txt(s, rx, kpi_y2 + Inches(0.95), Inches(5.8), Inches(0.18), "3. 투자단계 분포", sz=9, color=C_DARK, bold=True)
            stage_cnt = result_df["투자단계"].value_counts()
            for j, (stage, cnt) in enumerate(stage_cnt.items()):
                if j >= 4: break
                sx = rx + Inches(j * 1.45)
                _shape(s, sx, kpi_y2 + Inches(1.28), Inches(1.35), Inches(0.42), C_XPALE, C_PALE, radius=True)
                _txt(s, sx, kpi_y2 + Inches(1.30), Inches(1.35), Inches(0.16), stage, sz=9, color=C_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
                _txt(s, sx, kpi_y2 + Inches(1.47), Inches(1.35), Inches(0.18), f"{cnt}개사 ({cnt/n_cos*100:.0f}%)", sz=8, color=C_DARK, align=PP_ALIGN.CENTER)

        # ── 하단: 분석 과정(좌, 번호박스) + 인사이트(우) ──
        risk_bot = BODY_Y + Inches(2.9)
        risk_process = [
            f"HHI 산출: 투자비중 제곱합 × 10,000 = {hhi:,} → {hhi_label}",
            f"실현율 산출: 회수 {result_df['회수금액_백만원'].sum():,.0f}M / 전체 {total_val:,.0f}M = {realized_pct}%",
            f"MOIC<1.0x 스크리닝: {len(under)}개사 → {'모니터링 대상' if len(under) > 0 else '해당 없음'}",
        ]
        _process_list(s, M_LEFT, risk_bot, Inches(5.7), risk_process, box_h=Inches(0.42), gap=Inches(0.08), title="4. 리스크 분석 과정")

        risk_ins = [
            f"{'집중도 높음 — 특정 기업 부진 시 펀드 전체에 영향 가능' if hhi > 2500 else '적정 분산 — 개별 기업 리스크 제한적'}",
            f"{'현금 회수 초기 — 미실현 가치 현실화 모니터링 필요' if realized_pct < 30 else '회수 진행 중 — LP 배분 가능성 증가'}",
        ]
        _insight_panel(s, Inches(6.6), risk_bot + Inches(0.24), Inches(6.13), Inches(1.0), "RISK ASSESSMENT", risk_ins)
        slide_labels.append("리스크")

    # ════════════════════════════════════════════════
    # 8. J-Curve
    # ════════════════════════════════════════════════
    if _sec("J-Curve") and jcurve_df is not None and not jcurve_df.empty:
        s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
        mn = jcurve_df["누적현금흐름"].min()
        cr = jcurve_df["누적현금흐름"].iloc[-1]
        be = jcurve_df[jcurve_df["누적현금흐름"] >= 0]
        be_dt = str(be["날짜"].iloc[0])[:10] if not be.empty else "미도달"
        recovery = cr / abs(mn) * 100 if mn != 0 else 0
        _header(s, "J-CURVE ANALYSIS", "J-Curve 현금흐름 분석",
                f'최대 {abs(mn):,.0f}M 투자 후 현재 누적 CF {cr:,.0f}M으로, 회복률은 {recovery:.0f}%입니다.')

        # KPI 4개
        kw = Inches(2.85)
        _kpi_card(s, M_LEFT, BODY_Y, kw, Inches(0.90), "MAX DRAWDOWN", f"{mn:,.0f}M", "최대 누적 투자", C_RED)
        _kpi_card(s, M_LEFT + Inches(3.05), BODY_Y, kw, Inches(0.90), "현재 누적 CF", f"{cr:,.0f}M", "현재 순현금흐름", C_PRIMARY if cr >= 0 else C_RED)
        _kpi_card(s, M_LEFT + Inches(6.10), BODY_Y, kw, Inches(0.90), "회복률", f"{recovery:.0f}%", "현재CF / MaxDD", C_PRIMARY if recovery > 100 else C_ORANGE)
        _kpi_card(s, M_LEFT + Inches(9.15), BODY_Y, kw, Inches(0.90), "손익분기", be_dt, "BEP 달성 시점", C_PRIMARY if be_dt != "미도달" else C_GREY)

        # 차트 (좌측, 폭 좁혀 우측에 현금흐름 테이블 배치)
        chart_top = Inches(2.75)
        if include_charts:
            import plotly.graph_objects as go
            dates = jcurve_df["날짜"].astype(str).tolist()
            cf = jcurve_df["누적현금흐름"].tolist()
            fig = go.Figure()
            fig.add_trace(go.Scatter(x=dates, y=cf, mode="lines+markers",
                                     line=dict(color="#1b5e20", width=2.5),
                                     fill="tozeroy", fillcolor="rgba(27,94,32,0.08)",
                                     marker=dict(size=5, color=["#1b5e20" if v >= 0 else "#c62828" for v in cf])))
            fig.add_hline(y=0, line_dash="dash", line_color="#999")
            fig.update_layout(height=230, width=580, margin=dict(t=10, b=25, l=50, r=15),
                              paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)", showlegend=False,
                              xaxis=dict(showgrid=False), yaxis=dict(showgrid=True, gridcolor="#eee", title="백만원"))
            _chart_img(s, fig, 0.6, float(chart_top / 914400), 6.4, 2.6)

        # 우측: 현금흐름 테이블
        jc_hdrs = ["날짜", "현금흐름(M)", "누적(M)"]
        jc_cw = [Inches(1.3), Inches(1.05), Inches(1.05)]
        jc_rows_data = []
        for _, r in jcurve_df.tail(6).iterrows():
            jc_rows_data.append([str(r["날짜"])[:10], f'{r.get("현금흐름_백만원", 0):,.0f}', f'{r["누적현금흐름"]:,.0f}'])
        if jc_rows_data:
            _txt(s, Inches(7.4), chart_top, Inches(5), Inches(0.18), "현금흐름 상세", sz=8.5, color=C_GREY, bold=True)
            _table(s, Inches(7.4), chart_top + Inches(0.22), jc_hdrs, jc_rows_data, jc_cw, Inches(0.30))

        # 하단: 분석 과정(좌, 번호박스) + 인사이트(우)
        jc_stage = "투자 집행기" if cr < 0 else ("회수 전환기" if recovery < 100 else "수익 실현기")
        jc_bot = chart_top + Inches(2.70)
        jc_process = [
            f"{len(jcurve_df)}건의 현금흐름을 시계열로 누적 산출하여 J-Curve 형성",
            f"최대 누적 손실 {abs(mn):,.0f}M 시점을 투자 집행 피크로 확인",
            f"현재 단계는 {jc_stage}이며, 회복률은 {recovery:.0f}%",
        ]
        _process_list(s, M_LEFT, jc_bot, Inches(5.7), jc_process, box_h=Inches(0.38), gap=Inches(0.06))

        jc_ins = [
            f"손익분기: {'달성(' + be_dt + ')' if be_dt != '미도달' else '미도달 — 추가 회수 필요'}",
            f"{'회복률 100% 초과 — 원금 회수 완료' if recovery >= 100 else '원금 회수 진행 중 — 추가 모니터링 필요'}",
        ]
        _insight_panel(s, Inches(6.6), jc_bot + Inches(0.24), Inches(6.13), Inches(1.0), "J-CURVE INSIGHT", jc_ins)
        slide_labels.append("J-Curve")

    # ════════════════════════════════════════════════
    # 9. 시나리오 분석
    # ════════════════════════════════════════════════
    if _sec("시나리오") and scenario_df is not None and not scenario_df.empty:
        s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
        _header(s, "EXIT SCENARIO", f"회수 시나리오 — {scenario_company}",
                f'{scenario_company}의 Exit 배수별 예상 IRR을 시뮬레이션하여 최적 회수 타이밍을 분석합니다.')

        # 좌측: 차트
        if include_charts and "Exit 배수" in scenario_df.columns and "IRR (%)" in scenario_df.columns:
            import plotly.graph_objects as go
            fig = go.Figure(go.Bar(
                x=scenario_df["Exit 배수"].tolist(), y=scenario_df["IRR (%)"].tolist(),
                marker_color=["#1b5e20" if v >= 20 else "#43a047" if v >= 0 else "#c62828" for v in scenario_df["IRR (%)"]],
                text=[f"{v}%" for v in scenario_df["IRR (%)"]], textposition="outside", textfont=dict(size=9),
                marker_line_width=0,
            ))
            fig.add_hline(y=15, line_dash="dot", line_color="#999", annotation_text="Target 15%", annotation_font_size=8)
            fig.update_layout(height=280, width=600, margin=dict(t=10, b=30, l=40, r=15),
                              paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)", showlegend=False, bargap=0.3,
                              xaxis=dict(showgrid=False, title="Exit 배수"), yaxis=dict(showgrid=True, gridcolor="#eee", title="IRR (%)"))
            _chart_img(s, fig, 0.6, 1.50, 6.5, 3.5)

        # 우측: 시나리오 테이블
        sc_hdrs = ["Exit 배수", "IRR (%)", "회수금액(M)", "판정"]
        sc_cw = [Inches(0.8), Inches(0.8), Inches(1.2), Inches(0.8)]
        sc_rows = []
        for _, r in scenario_df.iterrows():
            irr_val = r.get("IRR (%)", 0)
            verdict = "달성" if irr_val >= 15 else "미달"
            exit_disp = r.get("Exit 배수", "")
            exit_disp = exit_disp if str(exit_disp).endswith(("x", "X")) else f"{exit_disp}x"
            sc_rows.append([exit_disp, f'{irr_val}%',
                            f'{r.get("회수금액 (백만원)", 0):,.0f}', verdict])
        _table(s, Inches(7.5), BODY_Y, sc_hdrs, sc_rows, sc_cw, Inches(0.28))

        # 하단 좌측: 분석 과정
        target_rows = scenario_df[scenario_df["IRR (%)"] >= 15] if "IRR (%)" in scenario_df.columns else pd.DataFrame()

        def _to_num(v):
            try:
                return float(str(v).rstrip("xX"))
            except (ValueError, TypeError):
                return None

        min_exit_val = None
        if not target_rows.empty:
            nums = [n for n in target_rows["Exit 배수"].apply(_to_num) if n is not None]
            min_exit_val = min(nums) if nums else None
        min_exit_disp = f"{min_exit_val}x" if min_exit_val is not None else "N/A"

        max_irr = scenario_df["IRR (%)"].max() if "IRR (%)" in scenario_df.columns else 0
        max_irr_exit = scenario_df.loc[scenario_df['IRR (%)'].idxmax(), 'Exit 배수'] if "IRR (%)" in scenario_df.columns else "-"
        max_irr_exit_disp = max_irr_exit if str(max_irr_exit).endswith(("x", "X")) else f"{max_irr_exit}x"

        sc_bot = Inches(5.3)
        sc_process = [
            f"{scenario_company}의 현재 투자금을 기준으로 {len(scenario_df)}개 Exit 배수 시나리오를 시뮬레이션",
            f"목표 IRR 15% 달성을 위해서는 최소 Exit {min_exit_disp} 이상이 필요",
            f"최대 IRR {max_irr}%는 Exit {max_irr_exit_disp} 시나리오에서 달성",
            f"Exit 배수가 1.0x 이하인 경우 원금 손실이 발생하므로 최소 1.5x 이상을 목표로 설정 권장",
        ]
        _process_list(s, M_LEFT, sc_bot, Inches(5.7), sc_process, box_h=Inches(0.38), gap=Inches(0.06))

        # 하단 우측: 인사이트
        min_exit_ok = min_exit_val is not None and min_exit_val <= 2.0
        sc_ins = [
            f"IRR 15% 달성 최소 배수: {min_exit_disp} — {'달성 가능성 높음' if min_exit_ok else '높은 Exit 배수 필요'}",
            f"LP 관점에서 Exit 타이밍과 배수 관리가 펀드 수익률을 결정짓는 핵심 변수입니다.",
        ]
        _insight_panel(s, Inches(6.6), sc_bot + Inches(0.24), Inches(6.13), Inches(1.4), "SCENARIO INSIGHT", sc_ins)
        slide_labels.append("시나리오")

    # ════════════════════════════════════════════════
    # 10. IRR Sensitivity
    # ════════════════════════════════════════════════
    if _sec("Sensitivity") and sensitivity_df is not None:
        s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
        _header(s, "IRR SENSITIVITY", f"IRR Sensitivity — {sensitivity_company}",
                f'{sensitivity_company}의 Exit 배수와 보유기간 조합에 따른 IRR 변동을 매트릭스로 분석합니다.')

        if include_charts:
            import plotly.graph_objects as go
            matrix = sensitivity_df.values.tolist()
            fig = go.Figure(data=go.Heatmap(
                z=matrix, x=list(sensitivity_df.columns), y=list(sensitivity_df.index),
                colorscale=[[0, "#c62828"], [0.15, "#e53935"], [0.3, "#ff9800"], [0.45, "#ffc107"],
                            [0.55, "#cddc39"], [0.7, "#66bb6a"], [0.85, "#43a047"], [1.0, "#1b5e20"]],
                zmid=15, text=[[f"{v}%" for v in row] for row in matrix],
                texttemplate="%{text}", textfont=dict(size=7, color="#fff"),
                showscale=False,
            ))
            fig.update_layout(height=300, width=750, margin=dict(t=10, b=30, l=50, r=15),
                              paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)",
                              xaxis_title="보유기간", yaxis_title="Exit 배수")
            _chart_img(s, fig, 0.6, 1.50, 8.0, 3.8)
        else:
            sens_hdrs = ["Exit\\기간"] + list(sensitivity_df.columns)[:8]
            sens_cw = [Inches(0.8)] + [Inches(0.8)] * min(8, len(sensitivity_df.columns))
            sens_rows = []
            for idx, row in sensitivity_df.iterrows():
                sens_rows.append([str(idx)] + [f"{row[c]}%" for c in list(sensitivity_df.columns)[:8]])
            _table(s, M_LEFT, BODY_Y, sens_hdrs, sens_rows, sens_cw, Inches(0.30))

        # 우측: 컬러 범례 + 해석 (차트 옆이 아닌 오른쪽 상단에 작게)
        _txt(s, Inches(9.0), BODY_Y, Inches(3.7), Inches(0.20), "SENSITIVITY 해석", sz=9, color=C_GREY, bold=True)
        sens_legend = [
            "■ 초록: IRR ≥15% (목표 달성)",
            "■ 노랑: IRR 5~15% (보통)",
            "■ 빨강: IRR <5% (저성과)",
        ]
        _multiline(s, Inches(9.0), BODY_Y + Inches(0.25), Inches(3.7), Inches(0.8), sens_legend, sz=9, color=C_DARK, spacing=5)

        # 하단: 분석 과정 + 인사이트 (차트 아래)
        sens_bot = Inches(5.5)
        sens_process = [
            f"① {sensitivity_company}의 현재 투자금 기준으로 Exit 배수(0.5~5.0x)와 보유기간(1~10년) 조합별 IRR을 산출",
            f"② 목표 IRR 15% 달성 가능 구간을 매트릭스에서 식별 — 초록 영역이 목표 달성 구간",
            f"③ LP 관점 최적 구간: Exit 배수 2.0x 이상 · 보유기간 3~5년이 IRR 극대화에 유리",
        ]
        _multiline(s, M_LEFT, sens_bot, Inches(6.0), Inches(1.0), sens_process, sz=9, color=C_DARK, spacing=5)

        sens_ins = [
            "단기 고배수 Exit이 IRR 극대화에 가장 유리합니다.",
            "장기 보유 시에도 2.0x 이상이면 안정적 수익이 가능합니다.",
        ]
        _insight_panel(s, Inches(6.8), sens_bot, Inches(5.9), Inches(0.85), "SENSITIVITY INSIGHT", sens_ins)
        slide_labels.append("Sensitivity")

    # ════════════════════════════════════════════════
    # 11. Waterfall 분배
    # ════════════════════════════════════════════════
    if _sec("Waterfall"):
        s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
        wf_inv = float(total_inv)
        wf_proc = float(result_df["현재가치_백만원"].sum() + result_df["회수금액_백만원"].sum())
        hurdle, carry, years = 8, 20, 5
        profit = max(0, wf_proc - wf_inv)
        hurdle_amt = wf_inv * ((1 + hurdle / 100) ** years - 1)
        rem = wf_proc
        s1 = min(rem, wf_inv); rem -= s1
        s2 = min(rem, hurdle_amt); rem -= s2
        gp_t = profit * carry / 100
        s3_gp = min(rem, gp_t); rem -= s3_gp
        s4_gp = rem * carry / 100; s4_lp = rem - s4_gp
        total_lp = s1 + s2 + s4_lp; total_gp = s3_gp + s4_gp
        lp_moic = total_lp / wf_inv if wf_inv > 0 else 0
        eff_carry = total_gp / profit * 100 if profit > 0 else 0

        _header(s, "WATERFALL DISTRIBUTION", "Waterfall 분배 시뮬레이션",
                f'LP에게 {total_lp:,.0f}M(MOIC {lp_moic:.2f}x), GP에게 {total_gp:,.0f}M이 분배되며 실효 Carry는 {eff_carry:.1f}%입니다.')

        # 좌측: 파라미터 + 차트
        _shape(s, M_LEFT, BODY_Y, Inches(6.3), Inches(0.35), C_XPALE, C_PALE, radius=True)
        _txt(s, M_LEFT + Inches(0.1), BODY_Y + Inches(0.06), Inches(6.1), Inches(0.2),
             f"투자금 {wf_inv:,.0f}M · 회수금 {wf_proc:,.0f}M · Hurdle {hurdle}% · Carry {carry}% · {years}년",
             sz=9, color=C_PRIMARY, bold=True, align=PP_ALIGN.CENTER)

        if include_charts:
            import plotly.graph_objects as go
            labels = ["① 원금반환", "② 우선수익", "③ GP Catch-up", "④ Carry Split"]
            fig = go.Figure()
            fig.add_trace(go.Bar(name="LP", x=labels, y=[s1, s2, 0, s4_lp],
                                 marker_color="#1b5e20",
                                 text=[f"{v:,.0f}" for v in [s1, s2, 0, s4_lp]],
                                 textposition="inside", textfont=dict(color="#fff", size=9)))
            fig.add_trace(go.Bar(name="GP", x=labels, y=[0, 0, s3_gp, s4_gp],
                                 marker_color="#a5d6a7",
                                 text=[f"{v:,.0f}" if v > 0 else "" for v in [0, 0, s3_gp, s4_gp]],
                                 textposition="inside", textfont=dict(size=9)))
            fig.update_layout(barmode="stack", height=250, width=550, margin=dict(t=10, b=30, l=40, r=10),
                              paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)",
                              legend=dict(orientation="h", y=-0.18), bargap=0.3,
                              yaxis=dict(showgrid=True, gridcolor="#eee", title="백만원"))
            _chart_img(s, fig, 0.6, float(BODY_Y / 914400) + 0.55, 6.0, 3.0)

        # 우측: 분배 상세 테이블 (파라미터 박스와 간격 확보)
        wf_hdrs = ["단계", "LP", "GP", "설명"]
        wf_cw = [Inches(1.2), Inches(1.0), Inches(1.0), Inches(2.0)]
        wf_rows = [
            ["① 원금반환", f"{s1:,.0f}", "-", "LP 투자 원금 반환"],
            ["② 우선수익", f"{s2:,.0f}", "-", f"Hurdle {hurdle}% × {years}년"],
            ["③ GP Catch-up", "-", f"{s3_gp:,.0f}", f"GP Carry 목표 보전"],
            ["④ Carry Split", f"{s4_lp:,.0f}", f"{s4_gp:,.0f}", f"LP {100-carry}% / GP {carry}%"],
            ["합계", f"{total_lp:,.0f}", f"{total_gp:,.0f}", f"총 {wf_proc:,.0f}M"],
        ]
        wf_tbl_x = Inches(7.4)
        wf_tbl_y = BODY_Y + Inches(0.55)
        _table(s, wf_tbl_x, wf_tbl_y, wf_hdrs, wf_rows, wf_cw, Inches(0.30))

        # 테이블 아래: 분배 구조 설명 (박스형 2x2)
        wf_desc_y = wf_tbl_y + Inches(0.30 * 6) + Inches(0.15)
        wf_step_items = [
            ("① 원금반환", "LP 투자 원금을 최우선 반환"),
            ("② 우선수익", f"Hurdle {hurdle}% × {years}년 우선배분"),
            ("③ GP Catch-up", "GP가 Carry 비율만큼 추가 수취"),
            ("④ Carry Split", f"잔여 LP{100-carry}% / GP{carry}% 분배"),
        ]
        wf_box_w = Inches(2.35); wf_box_h = Inches(0.46); wf_box_gap = Inches(0.10)
        for wi, (wt, wd) in enumerate(wf_step_items):
            wx = wf_tbl_x + (wi % 2) * (wf_box_w + wf_box_gap)
            wy = wf_desc_y + (wi // 2) * (wf_box_h + wf_box_gap)
            _shape(s, wx, wy, wf_box_w, wf_box_h, C_XPALE, C_PALE, radius=True)
            _txt(s, wx + Inches(0.10), wy + Inches(0.05), wf_box_w - Inches(0.20), Inches(0.18),
                 wt, sz=8.5, color=C_PRIMARY, bold=True)
            _txt(s, wx + Inches(0.10), wy + Inches(0.23), wf_box_w - Inches(0.20), Inches(0.20),
                 wd, sz=7.5, color=C_DARK)

        # 하단: KPI 3개 + 인사이트 — 4개 박스를 슬라이드 전폭에 균등 배치
        ky = Inches(5.65)
        kw = Inches(2.75)
        gap4 = Inches(0.15)
        ins_w = Inches(3.4)
        x0 = M_LEFT
        x1 = x0 + kw + gap4
        x2 = x1 + kw + gap4
        x3 = x2 + kw + gap4
        _kpi_card(s, x0, ky, kw, Inches(0.85), "LP 수취", f"{total_lp:,.0f}M", f"MOIC {lp_moic:.2f}x", C_PRIMARY)
        _kpi_card(s, x1, ky, kw, Inches(0.85), "GP Carry", f"{total_gp:,.0f}M", f"실효 {eff_carry:.1f}%", C_ACCENT)
        _kpi_card(s, x2, ky, kw, Inches(0.85), "수익 배분", f"{total_lp/wf_proc*100:.0f}% / {total_gp/wf_proc*100:.0f}%" if wf_proc > 0 else "-", "LP / GP", C_BLACK)

        wf_ins = [
            f"LP 수취 {total_lp:,.0f}M, MOIC {lp_moic:.2f}x 달성",
            f"{'LP 원금 이상 회수 완료' if lp_moic >= 1.0 else 'LP 원금 미달, 추가 회수 필요'}",
        ]
        _insight_panel(s, x3, ky, ins_w, Inches(1.15), "DISTRIBUTION", wf_ins)
        slide_labels.append("Waterfall")

    # ════════════════════════════════════════════════
    # 12. KVIC 시장 비교
    # ════════════════════════════════════════════════
    if _sec("KVIC") and kvic_sector_df is not None and not kvic_sector_df.empty:
        s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
        kvic_total = kvic_sector_df["총약정액(억원)"].sum()
        kvic_funds = int(kvic_sector_df["조합수"].sum())
        my_inv = result_df["투자금액_백만원"].sum() / 100
        avg_fund = kvic_total / kvic_funds if kvic_funds > 0 else 0

        _header(s, "KVIC MARKET COMPARISON", "KVIC 시장 비교",
                f'내 펀드 규모 {my_inv:,.0f}억원은 KVIC 전체 {kvic_total:,.0f}억원({kvic_funds:,}개 조합) 대비 평균 조합의 {my_inv/avg_fund:.1f}배 수준입니다.' if avg_fund > 0 else f'KVIC 전체 {kvic_total:,.0f}억원 시장과 내 펀드를 비교 분석합니다.')

        # KPI 행
        _kpi_card(s, M_LEFT, BODY_Y, Inches(2.85), Inches(0.90), "내 펀드 규모", f"{my_inv:,.0f}억",
                  f"KVIC의 {my_inv/kvic_total*100:.3f}%" if kvic_total > 0 else "", C_PRIMARY)
        _kpi_card(s, M_LEFT + Inches(3.05), BODY_Y, Inches(2.85), Inches(0.90), "KVIC 전체", f"{kvic_total:,.0f}억",
                  f"{kvic_funds:,}개 조합", C_ACCENT)
        _kpi_card(s, M_LEFT + Inches(6.10), BODY_Y, Inches(2.85), Inches(0.90), "평균 조합 대비",
                  f"{my_inv/avg_fund:.1f}배" if avg_fund > 0 else "-",
                  f"평균 {avg_fund:,.0f}억/조합", C_PRIMARY)
        _kpi_card(s, M_LEFT + Inches(9.15), BODY_Y, Inches(2.85), Inches(0.90), "투자 분야",
                  f"{len(kvic_sector_df)}개", "KVIC 섹터 수", C_BLACK)

        # 좌측: 차트
        if include_charts:
            import plotly.graph_objects as go
            top8 = kvic_sector_df.head(8)
            fig = go.Figure(go.Bar(
                x=top8["총약정액(억원)"].tolist(), y=top8["투자분야"].tolist(), orientation="h",
                marker_color="rgba(27,94,32,0.6)", marker_line_width=0,
                text=[f"{v:,.0f}억" for v in top8["총약정액(억원)"]], textposition="outside", textfont=dict(size=8),
            ))
            fig.update_layout(height=220, width=550, margin=dict(t=5, b=10, l=100, r=40),
                              paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)",
                              yaxis=dict(autorange="reversed", showgrid=False),
                              xaxis=dict(showgrid=True, gridcolor="#eee"), bargap=0.25)
            _chart_img(s, fig, 0.6, 2.55, 6.0, 2.8)

        # 우측: 섹터 테이블
        kv_hdrs = ["투자분야", "약정액(억)", "조합수", "비중(%)"]
        kv_cw = [Inches(1.5), Inches(0.9), Inches(0.7), Inches(0.7)]
        kv_rows = []
        for _, r in kvic_sector_df.head(8).iterrows():
            pct = r["총약정액(억원)"] / kvic_total * 100 if kvic_total > 0 else 0
            kv_rows.append([r["투자분야"], f'{r["총약정액(억원)"]:,.0f}', f'{int(r["조합수"])}', f'{pct:.1f}%'])
        _table(s, Inches(7.0), Inches(2.55), kv_hdrs, kv_rows, kv_cw, Inches(0.26))

        # 하단: 분석 과정(좌, 번호박스) + 인사이트(우)
        kv_bot = Inches(5.4)
        top_sector = kvic_sector_df.iloc[0]
        kv_process = [
            f"KVIC 공시 데이터 기준 {kvic_funds:,}개 조합, 총 {kvic_total:,.0f}억원 시장 규모 확인",
            f"내 펀드 {my_inv:,.0f}억원은 KVIC 평균 조합({avg_fund:,.0f}억) 대비 {my_inv/avg_fund:.1f}배" if avg_fund > 0 else "내 펀드 규모를 KVIC 평균 조합과 비교",
            f"KVIC 최대 분야: {top_sector['투자분야']} ({top_sector['총약정액(억원)']:,.0f}억, {top_sector['총약정액(억원)']/kvic_total*100:.0f}%)",
        ]
        _process_list(s, M_LEFT, kv_bot, Inches(5.7), kv_process, box_h=Inches(0.42), gap=Inches(0.08))

        kv_ins = [
            f"내 펀드는 KVIC 전체 시장의 {my_inv/kvic_total*100:.3f}%에 해당하며, 평균 조합 대비 {my_inv/avg_fund:.1f}배 규모입니다." if avg_fund > 0 else "",
            f"KVIC 시장 트렌드와 비교하여 포트폴리오 전략의 시장 적합성을 점검할 수 있습니다.",
        ]
        _insight_panel(s, Inches(6.6), kv_bot + Inches(0.24), Inches(6.13), Inches(1.0), "MARKET POSITIONING", [i for i in kv_ins if i])
        slide_labels.append("KVIC")

    # ════════════════════════════════════════════════
    # 13. 거시지표 · DART
    # ════════════════════════════════════════════════
    if _sec("거시") and (rate_df is not None or fx_df is not None):
        s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
        _header(s, "MACRO & DART", "거시지표 · DART 재무", "현재 금리 및 환율 환경을 분석하고, 포트폴리오 기업의 DART 재무제표를 점검합니다.")

        # 거시 KPI
        col_i = 0
        kw = Inches(2.85)
        if rate_df is not None and not rate_df.empty:
            latest_rate = rate_df["기준금리(%)"].iloc[-1]
            rate_chg = round(rate_df["기준금리(%)"].iloc[-1] - rate_df["기준금리(%)"].iloc[0], 2) if len(rate_df) > 1 else 0
            _kpi_card(s, M_LEFT + Inches(col_i * 3.05), BODY_Y, kw, Inches(0.90),
                      "기준금리", f"{latest_rate}%", f"변동 {rate_chg:+.2f}%p", C_PRIMARY)
            col_i += 1
        if fx_df is not None and not fx_df.empty:
            latest_fx = fx_df["원/달러(원)"].iloc[-1]
            fx_chg = round(fx_df["원/달러(원)"].iloc[-1] - fx_df["원/달러(원)"].iloc[0]) if len(fx_df) > 1 else 0
            _kpi_card(s, M_LEFT + Inches(col_i * 3.05), BODY_Y, kw, Inches(0.90),
                      "원/달러", f"{latest_fx:,.0f}원", f"변동 {fx_chg:+.0f}원", C_BLACK)
            col_i += 1
        if spread is not None:
            _kpi_card(s, M_LEFT + Inches(col_i * 3.05), BODY_Y, kw, Inches(0.90),
                      "IRR vs 금리", f"{spread:+.1f}%p", "펀드 스프레드",
                      C_PRIMARY if spread > 0 else C_RED)
            col_i += 1

        # DART 데이터 없을 때: 금리/환율 추이 차트로 중간 영역 채움
        _macro_chart_drawn = False
        if (dart_fin_df is None or dart_fin_df.empty) and include_charts:
            _macro_chart_drawn = True
            trend_y = BODY_Y + Inches(1.10)
            import plotly.graph_objects as go
            fig = go.Figure()
            if rate_df is not None and not rate_df.empty and "기준금리(%)" in rate_df.columns:
                fig.add_trace(go.Scatter(y=rate_df["기준금리(%)"].tolist(), mode="lines+markers",
                                         name="기준금리(%)", line=dict(color="#1b5e20", width=2.5),
                                         yaxis="y1"))
            if fx_df is not None and not fx_df.empty and "원/달러(원)" in fx_df.columns:
                fig.add_trace(go.Scatter(y=fx_df["원/달러(원)"].tolist(), mode="lines+markers",
                                         name="원/달러(원)", line=dict(color="#43a047", width=2.5, dash="dot"),
                                         yaxis="y2"))
            fig.update_layout(height=190, width=950, margin=dict(t=10, b=25, l=50, r=50),
                              paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)",
                              legend=dict(orientation="h", y=-0.20),
                              xaxis=dict(showgrid=False, title="기간"),
                              yaxis=dict(title="기준금리(%)", showgrid=True, gridcolor="#eee"),
                              yaxis2=dict(title="원/달러(원)", overlaying="y", side="right", showgrid=False))
            _chart_img(s, fig, 0.6, float(trend_y / 914400), 11.9, 2.1)

            # 분석 과정 박스 (차트 바로 아래, 겹침 방지 간격 확보)
            macro_proc_y = trend_y + Inches(2.35)
            macro_process = [
                f"기준금리·환율 데이터를 한국은행 ECOS API에서 조회기간 전체 시계열로 수집",
                f"펀드 평균 IRR과 기준금리를 비교하여 스프레드(초과수익) 산출",
                f"금리·환율 변동 추세를 포트폴리오 밸류에이션 및 해외투자 전략에 반영",
            ]
            _process_list(s, M_LEFT, macro_proc_y, Inches(11.9), macro_process,
                          box_h=Inches(0.36), gap=Inches(0.06), title="MACRO ANALYSIS PROCESS")

        # DART 재무 (있으면)
        if dart_fin_df is not None and not dart_fin_df.empty:
            dart_y = BODY_Y + Inches(1.15)
            _txt(s, M_LEFT, dart_y, Inches(5), Inches(0.2),
                 f"DART 재무제표 — {dart_company}", sz=11, color=C_BLACK, bold=True)

            # 차트
            if include_charts:
                import plotly.graph_objects as go
                cdf = dart_fin_df.copy()
                for col in ["매출액", "영업이익", "당기순이익"]:
                    if col in cdf.columns:
                        cdf[col] = cdf[col].apply(lambda x: round(x / 1e6) if pd.notna(x) else 0)
                fig = go.Figure()
                for col, clr in [("매출액", "#1b5e20"), ("영업이익", "#43a047"), ("당기순이익", "#a5d6a7")]:
                    if col in cdf.columns:
                        fig.add_trace(go.Bar(name=col, x=cdf["연도"].tolist(), y=cdf[col].tolist(),
                                             marker_color=clr, marker_line_width=0))
                fig.update_layout(barmode="group", height=220, width=500, margin=dict(t=5, b=25, l=40, r=10),
                                  paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)", bargap=0.3,
                                  legend=dict(orientation="h", y=-0.15), yaxis=dict(showgrid=True, gridcolor="#eee", title="백만원"))
                _chart_img(s, fig, 0.6, float(dart_y / 914400) + 0.30, 5.5, 2.8)

            # 우측: DART 테이블
            d_hdrs = ["연도", "매출액(M)", "영업이익(M)", "순이익(M)"]
            d_cw = [Inches(0.8), Inches(1.2), Inches(1.2), Inches(1.2)]
            d_rows = []
            for _, r in dart_fin_df.iterrows():
                yr = str(int(r["연도"])) if pd.notna(r.get("연도")) else ""
                rev = f"{r['매출액']/1e6:,.0f}" if pd.notna(r.get("매출액")) and r["매출액"] != 0 else "-"
                op = f"{r['영업이익']/1e6:,.0f}" if pd.notna(r.get("영업이익")) and r["영업이익"] != 0 else "-"
                ni = f"{r['당기순이익']/1e6:,.0f}" if pd.notna(r.get("당기순이익")) and r["당기순이익"] != 0 else "-"
                d_rows.append([yr, rev, op, ni])
            _table(s, Inches(6.8), BODY_Y + Inches(1.45), d_hdrs, d_rows, d_cw, Inches(0.28))

            # DART 재무 해석 (매출 성장률)
            try:
                valid = dart_fin_df[dart_fin_df["매출액"] > 0].sort_values("연도")
                if len(valid) >= 2:
                    rev_growth = (valid.iloc[-1]["매출액"] / valid.iloc[0]["매출액"]) ** (1 / max(len(valid) - 1, 1)) - 1
                    op_margin = valid.iloc[-1]["영업이익"] / valid.iloc[-1]["매출액"] * 100 if valid.iloc[-1]["매출액"] != 0 else 0
                    dart_note = f"{dart_company}: 연평균 매출 성장률 {rev_growth*100:.1f}%, 최근 영업이익률 {op_margin:.1f}%"
                else:
                    dart_note = f"{dart_company}: 재무 데이터 분석 중 — 추가 연도 데이터 확보 필요"
            except Exception:
                dart_note = ""
            if dart_note:
                _txt(s, Inches(6.8), BODY_Y + Inches(1.45) + Inches(0.28 * (len(d_rows) + 1)) + Inches(0.10),
                     Inches(5.2), Inches(0.4), dart_note, sz=8, color=C_GREY)

        # 하단: 거시경제 해석 (추이 차트+분석과정이 이미 그려졌으면 중복이라 생략)
        if not _macro_chart_drawn:
            macro_bot = Inches(5.3)
            macro_texts = []
            if rate_df is not None and not rate_df.empty:
                macro_texts.append(f"기준금리 {latest_rate}%: {'인하 기조로 PE/VC 투자환경 개선 기대' if rate_chg < 0 else '금리 인상기로 밸류에이션 부담 증가 가능'}" if 'latest_rate' in dir() else "")
            if fx_df is not None and not fx_df.empty:
                macro_texts.append(f"환율 {latest_fx:,.0f}원: {'원화 약세로 해외 투자 포트폴리오 환차익 가능성' if fx_chg > 0 else '원화 강세로 해외 자산 환산 시 불리'}" if 'latest_fx' in dir() else "")
            if spread is not None:
                macro_texts.append(f"펀드 스프레드 {spread:+.1f}%p: {'금리 대비 초과수익 확보' if spread > 0 else '금리 수준에 미달하는 수익률'}")
            macro_texts = [t for t in macro_texts if t]
            if macro_texts:
                _txt(s, M_LEFT, macro_bot, Inches(3), Inches(0.18), "MACRO IMPACT", sz=9, color=C_GREY, bold=True)
                _multiline(s, M_LEFT, macro_bot + Inches(0.20), C_WIDTH, Inches(1.0), macro_texts, sz=9, color=C_DARK, spacing=5)
        slide_labels.append("거시지표")

    # ════════════════════════════════════════════════
    # 14. AI 코멘터리
    # ════════════════════════════════════════════════
    if _sec("AI") and commentary:
        s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
        _header(s, "AI COMMENTARY", "AI 분석 코멘터리 — Claude",
                "Claude AI가 전체 포트폴리오 데이터를 종합 분석하여 자동 생성한 분기 코멘터리입니다.")

        # 코멘터리를 섹션 키워드 기준으로 분할 (헤더 줄과 본문을 합쳐 내용 손실 방지)
        section_titles = ["전체 성과 총평", "Top Performer 분석", "리스크 요인", "DPI vs RVPI 분석", "향후 전략 제언"]
        lines_all = [l.strip() for l in commentary.split("\n") if l.strip()]
        sections = {t: [] for t in section_titles}
        current = None
        for line in lines_all:
            matched = None
            for t in section_titles:
                # "1. 전체 성과 총평" 같은 헤더 줄(짧고 키워드를 포함) 인식
                if t in line and len(line) <= len(t) + 8:
                    matched = t
                    break
            if matched:
                current = matched
                continue
            if current:
                sections[current].append(line)

        box_titles, paras = [], []
        for t in section_titles:
            content = " ".join(sections[t]).strip()
            if content:
                box_titles.append(t)
                paras.append(content)

        if not paras:
            # 폴백: 섹션 키워드를 찾지 못하면 단락 단위로 분할
            paras = [p.strip() for p in commentary.split("\n\n") if p.strip()]
            if len(paras) <= 1:
                raw = " ".join(lines_all)
                sentences = [seg.strip() + "다." for seg in raw.split("다.") if seg.strip()]
                paras, chunk = [], []
                for sent in sentences:
                    chunk.append(sent)
                    if len(chunk) >= 2:
                        paras.append(" ".join(chunk)); chunk = []
                if chunk:
                    paras.append(" ".join(chunk))
            paras = paras[:5]
            box_titles = [f"코멘터리 {i+1}" for i in range(len(paras))]

        box_h = Inches(1.02) if len(paras) <= 5 else Inches(5.0 / max(len(paras), 1))
        box_gap = Inches(0.08)
        for i, (title, para) in enumerate(zip(box_titles, paras)):
            by = BODY_Y + i * (box_h + box_gap)
            _shape(s, M_LEFT, by, C_WIDTH, box_h, C_WHITE, C_PALE, radius=True)
            _shape(s, M_LEFT + Inches(0.12), by + Inches(0.10), Inches(0.26), Inches(0.26), C_PRIMARY, radius=True)
            _txt(s, M_LEFT + Inches(0.12), by + Inches(0.115), Inches(0.26), Inches(0.22),
                 str(i + 1), sz=10, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
            _txt(s, M_LEFT + Inches(0.48), by + Inches(0.08), Inches(4), Inches(0.22),
                 title, sz=10, color=C_PRIMARY, bold=True)
            _txt(s, M_LEFT + Inches(0.20), by + Inches(0.46), C_WIDTH - Inches(0.40), box_h - Inches(0.54),
                 para, sz=9, color=C_DARK)
        slide_labels.append("AI")

    # ════════════════════════════════════════════════
    # 마지막 본문: End of Document (Appendix 앞으로 이동)
    # ════════════════════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, C_PRIMARY)
    _shape(s, Inches(0), Inches(0), SW, SH, C_PRIMARY)
    _txt(s, Inches(0), Inches(2.8), SW, Inches(0.9),
         "End of Document", sz=48, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    _shape(s, Inches(1.5), Inches(3.8), Inches(10.333), Pt(2), C_LIGHT)
    _txt(s, Inches(0), Inches(4.1), SW, Inches(0.4),
         "이수빈", sz=16, color=C_LIGHT, align=PP_ALIGN.CENTER)
    _txt(s, Inches(0), Inches(4.6), SW, Inches(0.3),
         f"{fund_name}  ·  {quarter}  ·  PE/VC 분기 보고 도우미", sz=10, color=C_PALE, align=PP_ALIGN.CENTER)
    _txt(s, Inches(0), Inches(6.2), SW, Inches(0.3),
         "본 보고서는 자동 생성된 참고 자료이며, 투자 의사결정의 최종 근거로 사용할 수 없습니다.", sz=8, color=C_LIGHT, align=PP_ALIGN.CENTER)
    slide_labels.append("End")

    # ════════════════════════════════════════════════
    # Appendix 디바이더 슬라이드 — End of Document와 동일한 진한 초록 단색
    # ════════════════════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, C_PRIMARY)
    _shape(s, Inches(0), Inches(0), SW, SH, C_PRIMARY)  # 테마 배경 덮어쓰는 명시적 사각형
    _txt(s, Inches(0), Inches(2.8), SW, Inches(0.9),
         "Appendix", sz=48, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    _shape(s, Inches(1.5), Inches(3.8), Inches(10.333), Pt(2), C_LIGHT)
    _txt(s, Inches(0), Inches(4.1), SW, Inches(0.4),
         "주요 지표 계산식", sz=16, color=C_LIGHT, align=PP_ALIGN.CENTER)
    slide_labels.append("Appendix Divider")

    # ════════════════════════════════════════════════
    # Appendix p1: 핵심 성과 지표 계산식
    # ════════════════════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
    _header(s, "Appendix", "주요 지표 계산식 (1/2) — 성과 지표",
            "업계 표준 계산 산식을 기반으로, 투자 원금 대비 절대 수익(MOIC)부터 시간가치를 반영한 실질 수익률(IRR)까지 지표를 자동 산출합니다.")

    ex_row = result_df.iloc[0]
    ex_name = ex_row["회사명"]
    ex_inv = ex_row["투자금액_백만원"]
    ex_cur = ex_row["현재가치_백만원"]
    ex_rec = ex_row["회수금액_백만원"]
    ex_moic = ex_row["MOIC"]
    ex_dpi = ex_row["DPI"]
    ex_rvpi = ex_row["RVPI"]
    ex_tvpi = ex_row["TVPI"]
    ex_irr = ex_row["IRR(%)"]

    fb_w = Inches(3.85)
    fb_h = Inches(1.85)
    fb_gap = Inches(0.20)
    fb_y = BODY_Y + Inches(0.20)

    _formula_box(s, M_LEFT, fb_y, fb_w, fb_h, "MOIC (Multiple on Invested Capital)",
                 "MOIC = (현재가치 + 회수금액) / 투자원금",
                 f"EX. {ex_name}: MOIC = ({ex_cur:,.0f}+{ex_rec:,.0f}) / {ex_inv:,.0f} = {ex_moic}x",
                 "1.0x 미만=원금손실, 2.0x 이상=우수, 3.0x 이상=최상위")

    _formula_box(s, M_LEFT + fb_w + fb_gap, fb_y, fb_w, fb_h, "DPI (Distributions to Paid-In)",
                 "DPI = 회수금액 / 투자원금",
                 f"EX. {ex_name}: DPI = {ex_rec:,.0f} / {ex_inv:,.0f} = {ex_dpi}x",
                 "1.0x = 원금 회수 완료, 0 = 아직 미회수")

    _formula_box(s, M_LEFT + (fb_w + fb_gap) * 2, fb_y, fb_w, fb_h, "TVPI (Total Value to Paid-In)",
                 "TVPI = DPI + RVPI = (회수금액+현재가치) / 투자원금",
                 f"EX. {ex_name}: TVPI = {ex_dpi} + {ex_rvpi} = {ex_tvpi}x",
                 "MOIC와 유사하나 LP 출자금 기준, 2.0x 이상 = 우수")

    fb_y2 = fb_y + fb_h + Inches(0.25)
    _formula_box(s, M_LEFT, fb_y2, fb_w, fb_h, "IRR (Internal Rate of Return)",
                 "Σ CFt / (1+r)^t = 0 을 만족하는 r",
                 f"EX. {ex_name}: 투자 {ex_inv:,.0f}M → 회수 {ex_cur+ex_rec:,.0f}M → IRR = {ex_irr}%",
                 "10% 미만=저조, 15% 이상=우수, 25% 이상=최상위")

    _formula_box(s, M_LEFT + fb_w + fb_gap, fb_y2, fb_w, fb_h, "RVPI (Residual Value to Paid-In)",
                 "RVPI = 현재가치 / 투자원금",
                 f"EX. {ex_name}: RVPI = {ex_cur:,.0f} / {ex_inv:,.0f} = {ex_rvpi}x",
                 "펀드 초기 높음, 후기에도 높으면 Exit 지연 의심")

    # IRR vs MOIC 관계 박스
    rel_x = M_LEFT + (fb_w + fb_gap) * 2
    _shape(s, rel_x, fb_y2, fb_w, fb_h, C_BEIGE, C_BORDER, radius=True)
    _txt(s, rel_x + Inches(0.12), fb_y2 + Inches(0.08), fb_w - Inches(0.24), Inches(0.22),
         "IRR vs MOIC 관계", sz=11, color=C_BLACK, bold=True)
    rel_texts = [
        "시간가치의 차이: MOIC 2.0x를 3년에 달성하면 IRR 26%, 5년이면 IRR 15%로 절반 수준",
        "PE/VC 펀드의 일반적 목표: MOIC 2.0x 이상 + IRR 15% 이상",
        "빠른 Exit이 IRR에 유리하지만, 충분한 가치 성장 후 Exit이 MOIC에 유리 — 균형 필요",
    ]
    _multiline(s, rel_x + Inches(0.12), fb_y2 + Inches(0.34), fb_w - Inches(0.24), Inches(1.4),
               rel_texts, sz=8, color=C_DARK, spacing=4)

    slide_labels.append("Appendix p1")

    # ════════════════════════════════════════════════
    # Appendix p2: 고급 분석 지표 계산식 (HHI, IRR Sensitivity, Waterfall)
    # ════════════════════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
    _header(s, "Appendix", "주요 지표 계산식 (2/2) — 고급 분석 도구",
            "포트폴리오 집중도부터 회수 시뮬레이션까지, 본 보고서가 사용한 고급 분석 지표의 계산 로직을 안내합니다.")

    weights_ap = result_df["투자금액_백만원"] / result_df["투자금액_백만원"].sum()
    hhi_ap = round((weights_ap ** 2).sum() * 10000)

    fb2_w = Inches(3.85); fb2_h = Inches(1.85); fb2_gap = Inches(0.20)
    fb2_y = BODY_Y + Inches(0.20)

    _formula_box(s, M_LEFT, fb2_y, fb2_w, fb2_h, "HHI (Herfindahl-Hirschman Index)",
                 "HHI = Σ(투자비중)² × 10,000",
                 f"EX. 본 펀드: 투자비중 제곱합 × 10,000 = {hhi_ap:,}",
                 "1,500 미만=LOW, 1,500~2,500=MEDIUM, 2,500 이상=HIGH(집중)")

    _formula_box(s, M_LEFT + fb2_w + fb2_gap, fb2_y, fb2_w, fb2_h, "IRR Sensitivity Matrix",
                 "IRR = (Exit배수)^(1/보유기간) − 1",
                 "EX. Exit 2.0x, 보유 3년 → IRR = (2.0)^(1/3) − 1 = 26.0%",
                 "Exit 배수(행) × 보유기간(열) 조합을 매트릭스로 산출, 단기·고배수일수록 IRR 극대화")

    _formula_box(s, M_LEFT + (fb2_w + fb2_gap) * 2, fb2_y, fb2_w, fb2_h, "회수 시나리오 시뮬레이션",
                 "회수금액 = 투자원금 × Exit 배수",
                 "EX. 투자 1,000M × Exit 2.5x = 회수 2,500M → IRR 환산",
                 "Exit 배수 구간(0.5x~5.0x)별 IRR을 산출하여 목표 IRR 달성 최소 배수 도출")

    fb2_y2 = fb2_y + fb2_h + Inches(0.25)
    _formula_box(s, M_LEFT, fb2_y2, fb2_w, fb2_h, "Waterfall — ① 원금반환 · ② 우선수익",
                 "우선수익 = 투자원금 × [(1+Hurdle)^연수 − 1]",
                 "EX. 9,050M × [(1.08)^5 − 1] = 4,247M",
                 "LP가 투자 원금과 약정 Hurdle Rate(통상 8%) 수익을 최우선으로 수취")

    _formula_box(s, M_LEFT + fb2_w + fb2_gap, fb2_y2, fb2_w, fb2_h, "Waterfall — ③ GP Catch-up",
                 "GP Catch-up = MIN(잔여재원, 총수익×Carry%)",
                 "EX. 총수익 9,165M × 20% = 1,833M",
                 "GP가 전체 수익 대비 약정 Carry 비율(통상 20%)만큼 우선 배분받는 단계")

    _formula_box(s, M_LEFT + (fb2_w + fb2_gap) * 2, fb2_y2, fb2_w, fb2_h, "Waterfall — ④ Carry Split",
                 "LP/GP 잔여분배 = 잔여재원 × (1−Carry%) / Carry%",
                 "EX. 잔여 3,085M → LP 2,468M(80%) · GP 617M(20%)",
                 "①~③ 이후 남은 수익을 LP와 GP가 약정 비율(통상 80:20)로 최종 분배")

    slide_labels.append("Appendix p2")

    # ════════════════════════════════════════════════
    # Appendix p3: 시뮬레이터 계산 로직
    # ════════════════════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
    _header(s, "Appendix · simulator.py", "시뮬레이터 계산 로직",
            "PPTX 시나리오 슬라이드와 Streamlit Analysis 탭이 동일한 함수를 공유하여 화면의 숫자와 보고서 숫자가 항상 일치합니다.")

    fp3_w = Inches(3.85); fp3_h = Inches(1.85); fp3_gap = Inches(0.20)
    fp3_y = BODY_Y + Inches(0.20)

    # 실제 예시용 데이터
    _ex_p3 = result_df.iloc[0]
    _ex_inv_p3 = _ex_p3["투자금액_백만원"]
    _ex_irr_p3 = _ex_p3["IRR(%)"]

    _formula_box(s, M_LEFT, fp3_y, fp3_w, fp3_h,
                 "A-1. Exit 시뮬레이션 — simulate_exit()",
                 "회수금액 = 투자원금 × Exit배수",
                 f"EX. {_ex_inv_p3:,.0f}M × 2.5x = {_ex_inv_p3*2.5:,.0f}M\n"
                 f"IRR = XIRR([(투자일, −{_ex_inv_p3:,.0f}), (오늘, {_ex_inv_p3*2.5:,.0f})])",
                 "MOIC = DPI = Exit배수(전액현금화가정). IRR은 날짜 기반 XIRR — 단순 (배수)^(1/년)−1의 근사치가 아님")

    _formula_box(s, M_LEFT + fp3_w + fp3_gap, fp3_y, fp3_w, fp3_h,
                 "A-2. 목표 IRR 역산 — optimal_exit_timing()",
                 "필요배수 = (1 + 목표IRR%)^경과연수",
                 "EX. 목표 IRR 20%, 투자 후 3년 경과\n"
                 "→ (1.20)^3 = 1.728x → 최소 MOIC 1.73x 필요",
                 "현재MOIC ≥ 필요배수 → 지금 Exit해도 목표 IRR 달성\n미달 → 보유 연장 또는 가치 제고 후 Exit 권장")

    # XIRR 수치해석 설명 박스 (formula_box 대신 커스텀)
    xirr_x = M_LEFT + (fp3_w + fp3_gap) * 2
    _shape(s, xirr_x, fp3_y, fp3_w, fp3_h, C_WHITE, C_BORDER, radius=True)
    _txt(s, xirr_x + Inches(0.14), fp3_y + Inches(0.10), fp3_w - Inches(0.28), Inches(0.24),
         "XIRR 수치해석 방식 (scipy.optimize.brentq)", sz=9.5, color=C_BLACK, bold=True)
    _shape(s, xirr_x + Inches(0.14), fp3_y + Inches(0.40), fp3_w - Inches(0.28), Inches(0.34), C_PALE, radius=True)
    _txt(s, xirr_x + Inches(0.20), fp3_y + Inches(0.44), fp3_w - Inches(0.40), Inches(0.28),
         "Σ CFt / (1+r)^t = 0  →  r을 brentq로 수치해석", sz=9, color=C_PRIMARY, bold=True)
    _multiline(s, xirr_x + Inches(0.14), fp3_y + Inches(0.84), fp3_w - Inches(0.28), Inches(0.52),
               ["· 실제 날짜(일 수) 기반 연환산 — 정밀 계산",
                "· 비정수 보유기간(예: 2.5년)에서 근사치와 오차 발생",
                "· irr.py의 xirr() 함수 통일 — 화면·보고서 숫자 일치"],
               sz=8, color=C_DARK, spacing=4)
    crit_y3 = fp3_y + Inches(1.46)
    _shape(s, xirr_x + Inches(0.14), crit_y3, fp3_w - Inches(0.28), Inches(0.30), C_BEIGE, radius=True)
    _txt(s, xirr_x + Inches(0.20), crit_y3 + Inches(0.07), fp3_w - Inches(0.40), Inches(0.20),
         "해석기준: 비정수 보유기간일수록 근사치 오차 증가 → XIRR이 필수", sz=7.5, color=C_DARK)

    # 하단 Row: 데이터 일관성 + 활용 시점
    fp3_y2 = fp3_y + fp3_h + Inches(0.25)

    # 데이터 일관성 박스
    _shape(s, M_LEFT, fp3_y2, fp3_w, fp3_h, C_XPALE, C_BORDER, radius=True)
    _txt(s, M_LEFT + Inches(0.14), fp3_y2 + Inches(0.10), fp3_w - Inches(0.28), Inches(0.22),
         "데이터 일관성 보장", sz=10.5, color=C_BLACK, bold=True)
    _shape(s, M_LEFT + Inches(0.14), fp3_y2 + Inches(0.38), fp3_w - Inches(0.28), Pt(2), C_PRIMARY)
    _process_list(s, M_LEFT + Inches(0.14), fp3_y2 + Inches(0.46), fp3_w - Inches(0.28),
                  ["PPTX 시나리오 슬라이드 (app.py ~1667행)",
                   "Streamlit Analysis 탭 (app.py ~1078행)",
                   "→ 동일한 simulate_exit() 호출 → 숫자 항상 일치"],
                  box_h=Inches(0.38), gap=Inches(0.06), title="공유 함수 호출 경로")

    # 시뮬레이터 활용 시점 박스
    _insight_panel(s, M_LEFT + fp3_w + fp3_gap, fp3_y2, fp3_w * 2 + fp3_gap, fp3_h,
                   "시뮬레이터 활용 시점 & 해석 가이드",
                   ["Exit 배수별 IRR 범위 확인 → 목표 IRR 달성 최소 배수 도출",
                    "목표 IRR 역산 → '지금 팔면 목표 수익률 나오나?' 즉시 판단",
                    "IRR Sensitivity 매트릭스와 병행 → Exit 타이밍 전략 수립",
                    "전체 포트폴리오 Exit 시나리오 비교 → LP 기대수익 커뮤니케이션",
                    "주의: IRR Sensitivity는 단순 근사치, 회수 시나리오는 XIRR 정밀 계산"])

    slide_labels.append("Appendix p3")

    # ════════════════════════════════════════════════
    # Appendix p4: 밸류에이션 자동 산출 로직
    # ════════════════════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
    # 제목 짧게 — "valuation_fetcher.py"는 섹션 라벨로 이동
    _header(s, "Appendix · valuation_fetcher.py", "밸류에이션 자동 산출 로직",
            "상장사는 네이버 금융 시총, 비상장사는 DART 재무 × 섹터 P/S 배수로 현재가치를 자동 산출합니다 (32개 섹터 정의).")

    fp4_y = BODY_Y + Inches(0.15)

    # 좌측: 조회 우선순위 (3단계 Fallback) — 직접 배치
    _fp4_bh = Inches(0.38)   # 아이템 박스 높이
    _fp4_gap = Inches(0.07)  # 아이템 간격
    _fp4_total = Inches(0.42) + 3 * (_fp4_bh + _fp4_gap) + Inches(0.26)  # 컨테이너 전체 높이
    _shape(s, M_LEFT, fp4_y, Inches(3.85), _fp4_total, C_WHITE, C_BORDER, radius=True)
    _txt(s, M_LEFT + Inches(0.14), fp4_y + Inches(0.10), Inches(3.57), Inches(0.20),
         "조회 우선순위 — 3단계 Fallback", sz=10.5, color=C_BLACK, bold=True)
    _shape(s, M_LEFT + Inches(0.14), fp4_y + Inches(0.36), Inches(3.57), Pt(2), C_PRIMARY)
    _fb_items = [
        ("①", "상장사 (최우선)",   "네이버 금융 시총 × 지분율"),
        ("②", "비상장사 / ①실패", "DART 매출 × P/S배수 × 지분율"),
        ("③", "둘 다 실패",       "None → 수동 입력값 유지"),
    ]
    for _fi, (_fn, _ft, _fd) in enumerate(_fb_items):
        _fy = fp4_y + Inches(0.46) + _fi * (_fp4_bh + _fp4_gap)
        _shape(s, M_LEFT + Inches(0.14), _fy, Inches(3.57), _fp4_bh, C_XPALE, C_BORDER, radius=True)
        _shape(s, M_LEFT + Inches(0.18), _fy + Inches(0.05), Inches(0.26), Inches(0.26), C_PRIMARY, radius=True)
        _txt(s, M_LEFT + Inches(0.18), _fy + Inches(0.07), Inches(0.26), Inches(0.22),
             _fn, sz=9, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
        _txt(s, M_LEFT + Inches(0.52), _fy + Inches(0.05), Inches(3.0), Inches(0.15),
             _ft, sz=8.5, color=C_BLACK, bold=True)
        _txt(s, M_LEFT + Inches(0.52), _fy + Inches(0.21), Inches(3.0), Inches(0.14),
             _fd, sz=8, color=C_DARK)
    _parallel_y = fp4_y + Inches(0.46) + 3 * (_fp4_bh + _fp4_gap) + Inches(0.04)
    _txt(s, M_LEFT + Inches(0.14), _parallel_y, Inches(3.57), Inches(0.18),
         "병렬 처리: ThreadPoolExecutor(max_workers=6) — 8개사 기준 약 5~10초",
         sz=7, color=C_GREY)

    # 중앙: 상장사 방식 — 수식 줄여서 한 줄에 맞춤
    _formula_box(s, M_LEFT + Inches(4.05), fp4_y, Inches(3.85), _fp4_total,
                 "B-1. 상장사 — 네이버 금융 시총",
                 "현재가치(M) = 시총(억) × 100 × 지분율%",
                 "EX. 시총 5,000억, 지분율 3%\n→ 5,000 × 100 × 0.03 = 15,000백만원",
                 "네이버 금융 HTML 파싱(BeautifulSoup) — 구조 변경 시 ②로 자동 fallback")

    # 우측: 비상장사 방식 — 높이 통일
    _formula_box(s, M_LEFT + Inches(8.10), fp4_y, Inches(3.85), _fp4_total,
                 "B-2. 비상장사 — 섹터 멀티플 3방법 평균",
                 "평균가치 = (P/S + EV/EBITDA? + P/E?) / 방법수",
                 "EX. 바이오(P/S 8.0x), 매출 100억, 영업이익 20억, 순이익 10억\n"
                 "P/S=800억 / EV/EBITDA=240억 / P/E=240억 → 평균 426.7억",
                 "32개 섹터 P/S배수 정의(미정의=3.0x) / EV/EBITDA=P/S×1.5 / P/E=MAX(P/S×3, 10)")

    # 하단: 섹터 P/S 배수 샘플 테이블 + 투명성 설명
    fp4_y2 = fp4_y + _fp4_total + Inches(0.15)   # 상단 박스 아래 여백

    # 주요 섹터 P/S 배수 샘플 (일부)
    sec_sample = [
        ("바이오/헬스케어", "8.0x"), ("의료AI", "7.0x"), ("AI/딥테크", "7.0x"),
        ("반도체/하드웨어", "6.0x"), ("SaaS/클라우드", "5.0x"), ("핀테크", "4.5x"),
        ("모빌리티", "3.0x"), ("이커머스", "2.5x"), ("부동산/건설", "2.0x"), ("미정의", "3.0x"),
    ]
    _txt(s, M_LEFT, fp4_y2 + Inches(0.08), Inches(6), Inches(0.18),
         "주요 섹터 P/S 배수 (KOSDAQ 동종업계 중앙값 기반, 총 32개 섹터 정의)", sz=9, color=C_GREY, bold=True)
    col_w4 = Inches(1.10); row_h4 = Inches(0.25)
    hdrs4  = ["섹터", "P/S배수"] * 5
    cw4    = [Inches(1.10), Inches(0.55)] * 5
    rows4  = []
    for k in range(0, len(sec_sample), 2):
        pair = sec_sample[k:k+2]
        row_data = []
        for nm, mult in pair:
            row_data += [nm, mult]
        while len(row_data) < 10:
            row_data += ["", ""]
        rows4.append(row_data)
    _table(s, M_LEFT, fp4_y2 + Inches(0.30), hdrs4, rows4, cw4, row_h4)

    # 투명성 장치 설명
    _shape(s, M_LEFT, fp4_y2 + Inches(1.60), C_WIDTH, Inches(0.70), C_XPALE, C_BORDER, radius=True)
    _shape(s, M_LEFT, fp4_y2 + Inches(1.60), Inches(0.06), Inches(0.70), C_PRIMARY)
    _txt(s, M_LEFT + Inches(0.18), fp4_y2 + Inches(1.68), Inches(3.0), Inches(0.18),
         "투명성 장치", sz=9, color=C_PRIMARY, bold=True)
    _txt(s, M_LEFT + Inches(0.18), fp4_y2 + Inches(1.88), C_WIDTH - Inches(0.34), Inches(0.36),
         "모든 계산 결과에 근거 문자열이 함께 저장됩니다 — 예) '2024년 기준 | 매출 100.0억 | 영업이익률 20.0% | "
         "[P/S=800억 / EV/EBITDA=240억 / P/E=240억] × 지분 5%' → 클릭 한 번으로 산출 근거 전체 확인 가능",
         sz=8, color=C_DARK)

    slide_labels.append("Appendix p4")

    # ── 페이지 번호 ─────────────────────────────────
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides):
        if 0 < i < total - 1:
            _page_num(slide, i + 1, total)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
